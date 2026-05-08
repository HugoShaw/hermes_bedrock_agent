#!/usr/bin/env python3
"""
Phase 1: Download S3 files, extract text, chunk documents, extract entities and relations with provenance.
"""
import boto3
import json
import os
import re
import hashlib
import time
from pathlib import Path
from collections import defaultdict
from datetime import datetime, timezone

# === Configuration ===
BUCKET = "s3-hulftchina-rd"
PREFIX = "Murata/"
REGION = "ap-northeast-1"
LOCAL_ROOT = Path.home() / "projects/data/Murata"
EXTRACTED_DIR = Path.home() / "projects/data/extracted"
OUTPUT_DIR = Path.home() / "projects/data/output"
MAX_TEXT = 200_000  # Cap per-file text for regex processing
CHUNK_SIZE = 1500   # chars per chunk
CHUNK_OVERLAP = 200 # overlap between chunks

# Ensure dirs exist
LOCAL_ROOT.mkdir(parents=True, exist_ok=True)
EXTRACTED_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# === Helper Functions ===
def node_id(text: str) -> str:
    return "n_" + hashlib.md5(text.encode()).hexdigest()[:12]

def edge_id(src: str, pred: str, tgt: str) -> str:
    return "e_" + hashlib.md5(f"{src}:{pred}:{tgt}".encode()).hexdigest()[:12]

def chunk_id(file_path: str, idx: int) -> str:
    return "chunk_" + hashlib.md5(f"{file_path}:{idx}".encode()).hexdigest()[:10]

def safe_text(t, mx=500):
    if not isinstance(t, str): t = str(t)
    t = t.replace("\\", "\\\\").replace("'", "\\'").replace('"', '\\"')
    return re.sub(r'[\n\r\t]+', ' ', t).strip()[:mx]

def extract_plain(path):
    for enc in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return "[binary or unknown encoding - skipped]"

def extract_docx(path):
    try:
        from docx import Document
        doc = Document(str(path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[DOCX extraction error: {e}]"

def extract_xlsx(path):
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(max_row=200, values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    parts.append(row_text)
        wb.close()
        return "\n".join(parts)
    except Exception as e:
        return f"[XLSX extraction error: {e}]"

def extract_xls(path):
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        parts = []
        for sheet in wb.sheets():
            parts.append(f"=== Sheet: {sheet.name} ===")
            for rx in range(min(sheet.nrows, 200)):
                row_text = " | ".join(str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols) if sheet.cell_value(rx, cx))
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[XLS extraction error: {e}]"

def extract_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            parts.append(para.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[PPTX extraction error: {e}]"

def extract_pdf(path):
    try:
        from pdfminer.high_level import extract_text
        return extract_text(str(path))
    except Exception as e:
        return f"[PDF extraction error: {e}]"

# === Step 1: Download files from S3 ===
def download_s3_files():
    print("=" * 60)
    print("STEP 1: Downloading files from S3...")
    print("=" * 60)
    s3 = boto3.client("s3", region_name=REGION)
    paginator = s3.get_paginator("list_objects_v2")
    
    downloaded = []
    skipped_ext = {'.png', '.jpg', '.gif', '.db', '.jar', '.DS_Store', '.css'}
    
    for page in paginator.paginate(Bucket=BUCKET, Prefix=PREFIX):
        for obj in page.get("Contents", []):
            key = obj["Key"]
            size = obj["Size"]
            
            # Skip directories, images, binaries
            if key.endswith('/') or size == 0:
                continue
            ext = os.path.splitext(key)[1].lower()
            basename = os.path.basename(key)
            if ext in skipped_ext or basename.startswith('.') or basename.startswith('~$'):
                continue
            
            rel_path = key.replace(PREFIX, "", 1)
            local_path = LOCAL_ROOT / rel_path
            local_path.parent.mkdir(parents=True, exist_ok=True)
            
            if not local_path.exists() or local_path.stat().st_size != size:
                s3.download_file(BUCKET, key, str(local_path))
            
            downloaded.append({
                "key": key,
                "s3_path": f"s3://{BUCKET}/{key}",
                "local_path": str(local_path),
                "relative_path": rel_path,
                "size": size,
                "extension": ext,
                "file_name": basename,
                "modified": obj["LastModified"].isoformat()
            })
    
    print(f"  Downloaded/verified: {len(downloaded)} files")
    return downloaded

# === Step 2: Extract text from files ===
def extract_texts(files):
    print("\n" + "=" * 60)
    print("STEP 2: Extracting text from files...")
    print("=" * 60)
    
    manifest = []
    extractor_map = {
        '.java': 'plain', '.sql': 'plain', '.xml': 'plain',
        '.txt': 'plain', '.properties': 'plain', '.jsp': 'plain',
        '.md': 'plain', '.mf': 'plain', '.csv': 'plain',
        '.iml': 'plain',
        '.docx': 'docx', '.xlsx': 'xlsx', '.xls': 'xls',
        '.pptx': 'pptx', '.pdf': 'pdf'
    }
    
    for f in files:
        ext = f["extension"]
        local_path = Path(f["local_path"])
        extractor = extractor_map.get(ext, None)
        
        if extractor is None:
            continue
        
        try:
            if extractor == 'plain':
                text = extract_plain(local_path)
            elif extractor == 'docx':
                text = extract_docx(local_path)
            elif extractor == 'xlsx':
                text = extract_xlsx(local_path)
            elif extractor == 'xls':
                text = extract_xls(local_path)
            elif extractor == 'pptx':
                text = extract_pptx(local_path)
            elif extractor == 'pdf':
                text = extract_pdf(local_path)
            else:
                text = ""
        except Exception as e:
            text = f"[Error: {e}]"
        
        # Save extracted text
        safe_name = f["relative_path"].replace("/", "__").replace(" ", "_")
        extracted_path = EXTRACTED_DIR / f"{safe_name}.txt"
        extracted_path.write_text(text, encoding="utf-8")
        
        manifest.append({
            **f,
            "extracted_to": str(extracted_path),
            "extractor": extractor,
            "char_count": len(text),
            "is_empty": len(text.strip()) < 10
        })
    
    # Save manifest
    manifest_path = OUTPUT_DIR / "manifest.json"
    with open(manifest_path, 'w', encoding='utf-8') as fp:
        json.dump({"files": manifest, "total": len(manifest)}, fp, ensure_ascii=False, indent=2)
    
    non_empty = [m for m in manifest if not m["is_empty"]]
    print(f"  Extracted: {len(manifest)} files, {len(non_empty)} non-empty")
    return manifest

# === Step 3: Chunk documents ===
def chunk_documents(manifest):
    print("\n" + "=" * 60)
    print("STEP 3: Chunking documents...")
    print("=" * 60)
    
    all_chunks = []
    for entry in manifest:
        if entry["is_empty"]:
            continue
        
        text = Path(entry["extracted_to"]).read_text(encoding="utf-8")
        if len(text) > MAX_TEXT:
            text = text[:MAX_TEXT]
        
        # Semantic chunking by paragraphs/sections first
        chunks = []
        current_chunk = ""
        paragraphs = text.split('\n')
        
        for para in paragraphs:
            if len(current_chunk) + len(para) + 1 > CHUNK_SIZE:
                if current_chunk.strip():
                    chunks.append(current_chunk.strip())
                # Start new chunk with overlap
                overlap_start = max(0, len(current_chunk) - CHUNK_OVERLAP)
                current_chunk = current_chunk[overlap_start:] + "\n" + para
            else:
                current_chunk += "\n" + para
        
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        # If no chunks produced, make one
        if not chunks:
            chunks = [text[:CHUNK_SIZE]]
        
        for idx, chunk_text in enumerate(chunks):
            cid = chunk_id(entry["relative_path"], idx)
            all_chunks.append({
                "chunk_id": cid,
                "source_path": entry["s3_path"],
                "source_file": entry["file_name"],
                "relative_path": entry["relative_path"],
                "source_type": classify_doc_type(entry),
                "chunk_index": idx,
                "total_chunks": len(chunks),
                "text": chunk_text,
                "char_count": len(chunk_text)
            })
    
    # Save chunks
    chunks_path = OUTPUT_DIR / "chunks.json"
    with open(chunks_path, 'w', encoding='utf-8') as fp:
        json.dump(all_chunks, fp, ensure_ascii=False, indent=2)
    
    print(f"  Generated: {len(all_chunks)} chunks from {len([m for m in manifest if not m['is_empty']])} files")
    return all_chunks

def classify_doc_type(entry):
    """Classify document type for provenance."""
    path = entry["relative_path"].lower()
    ext = entry["extension"]
    
    if ext in ('.sql',):
        return "ddl"
    elif ext in ('.java',):
        return "code"
    elif ext in ('.xml', '.properties'):
        return "config"
    elif ext in ('.jsp',):
        return "code"
    elif ext in ('.docx',):
        if '手册' in path or 'manual' in path.lower():
            return "user_manual"
        elif '需求' in path or '设计' in path:
            return "design_doc"
        return "design_doc"
    elif ext in ('.xlsx', '.xls'):
        if 'qa' in path.lower() or 'テスト' in path:
            return "test"
        return "design_doc"
    elif ext in ('.pptx',):
        return "design_doc"
    elif ext in ('.csv',):
        return "data"
    elif ext in ('.txt',):
        if 'sql' in path.lower() or 'insert' in path.lower():
            return "sql"
        return "design_doc"
    return "other"

# === Step 4: Entity Extraction with Provenance ===
# Regex patterns for Oracle DDL
RE_CREATE_TABLE = re.compile(
    r'CREATE\s+TABLE\s+(?:"?\w+"?\.)?\"?(\w+)\"?\s*\(', re.I)
RE_CREATE_VIEW = re.compile(
    r'CREATE\s+(?:OR\s+REPLACE\s+)?(?:FORCE\s+)?VIEW\s+(?:"?\w+"?\.)?\"?(\w+)\"?\s', re.I)
RE_COMMENT_COL = re.compile(
    r"COMMENT\s+ON\s+COLUMN\s+\"?\w+\"?\.\"?(\w+)\"?\.\"?(\w+)\"?\s+IS\s+'([^']*)'", re.I)
RE_COMMENT_TBL = re.compile(
    r"COMMENT\s+ON\s+TABLE\s+\"?\w+\"?\.\"?(\w+)\"?\s+IS\s+'([^']*)'", re.I)
RE_JAVA_CLASS = re.compile(r'public\s+(?:abstract\s+)?class\s+(\w+)', re.I)
RE_JAVA_INTERFACE = re.compile(r'public\s+interface\s+(\w+)', re.I)
RE_AUTOWIRED = re.compile(r'@Autowired\s+(?:private\s+)?(\w+)\s+\w+', re.I)
RE_IMPORT = re.compile(r'import\s+com\.hulftchina\.[\w.]*\.(\w+);')
RE_JSP_ACTION = re.compile(r'action=["\']([^"\']+)["\']', re.I)
RE_STRUTS_ACTION = re.compile(r'<action\s+name="([^"]+)"[^>]*class="([^"]+)"', re.I)
RE_INSERT_FROM = re.compile(r'INSERT\s+INTO\s+\"?(\w+)\"?\s.*?SELECT\s.*?FROM\s+\"?(\w+)\"?', re.I | re.S)
RE_STATUS_PATTERN = re.compile(r"STATUS.*?(\d+[：:].+?)(?:'|$)", re.I)
RE_SYSTEM_REF = re.compile(r'(iMaps|SUN|HDS|HULFT|Oracle|Murata|ERP|MDW)', re.I)

def ctx(text, pos, window=200):
    """Extract context window around a regex match position."""
    return re.sub(r'[\n\r\t]+', ' ', text[max(0, pos-window):min(len(text), pos+window)].strip())

def extract_entities_and_relations(chunks, manifest):
    print("\n" + "=" * 60)
    print("STEP 4: Extracting entities and relations with provenance...")
    print("=" * 60)
    
    entities = {}  # id -> entity dict
    relations = []  # list of relation dicts
    
    # Build chunk lookup by file
    chunks_by_file = defaultdict(list)
    for c in chunks:
        chunks_by_file[c["relative_path"]].append(c)
    
    # Track tables, views, classes for cross-referencing
    tables = {}  # name -> node_id
    views = {}
    classes = {}
    interfaces = {}
    systems = set()
    
    for entry in manifest:
        if entry["is_empty"]:
            continue
        
        text = Path(entry["extracted_to"]).read_text(encoding="utf-8")
        if len(text) > MAX_TEXT:
            text = text[:MAX_TEXT]
        
        file_name = entry["file_name"]
        rel_path = entry["relative_path"]
        s3_path = entry["s3_path"]
        source_type = classify_doc_type(entry)
        file_chunks = chunks_by_file.get(rel_path, [])
        
        # Find best matching chunk for a position
        def find_chunk(pos):
            char_count = 0
            for c in file_chunks:
                char_count += c["char_count"]
                if pos < char_count:
                    return c["chunk_id"]
            return file_chunks[0]["chunk_id"] if file_chunks else "chunk_unknown"
        
        # === Extract Tables ===
        for m in RE_CREATE_TABLE.finditer(text):
            tname = m.group(1).upper()
            nid = node_id(f"Table:{tname}")
            evidence_text = ctx(text, m.start(), 400)
            cid = find_chunk(m.start())
            
            if nid not in entities:
                entities[nid] = {
                    "id": nid,
                    "label": tname,
                    "type": "Table",
                    "canonical_name": tname,
                    "aliases": [tname],
                    "provenance": []
                }
            entities[nid]["provenance"].append({
                "source_type": source_type,
                "source_path": s3_path,
                "source_chunk_id": cid,
                "source_text": evidence_text[:500],
                "confidence": 0.98
            })
            tables[tname] = nid
        
        # === Table comments ===
        for m in RE_COMMENT_TBL.finditer(text):
            tname = m.group(1).upper()
            comment = m.group(2)
            nid = tables.get(tname)
            if nid and nid in entities:
                entities[nid].setdefault("description", comment)
        
        # === Extract Views ===
        for m in RE_CREATE_VIEW.finditer(text):
            vname = m.group(1).upper()
            nid = node_id(f"View:{vname}")
            evidence_text = ctx(text, m.start(), 400)
            cid = find_chunk(m.start())
            
            if nid not in entities:
                entities[nid] = {
                    "id": nid,
                    "label": vname,
                    "type": "View",
                    "canonical_name": vname,
                    "aliases": [vname],
                    "provenance": []
                }
            entities[nid]["provenance"].append({
                "source_type": source_type,
                "source_path": s3_path,
                "source_chunk_id": cid,
                "source_text": evidence_text[:500],
                "confidence": 0.95
            })
            views[vname] = nid
        
        # === Extract Java Classes ===
        if entry["extension"] == ".java":
            for m in RE_JAVA_CLASS.finditer(text):
                cname = m.group(1)
                nid = node_id(f"JavaClass:{cname}")
                evidence_text = ctx(text, m.start(), 300)
                cid = find_chunk(m.start())
                
                # Determine subtype
                subtype = "JavaClass"
                if "Action" in cname:
                    subtype = "Controller"
                elif "Service" in cname:
                    subtype = "Service"
                elif "Dao" in cname:
                    subtype = "DataAccess"
                elif "Entity" in cname or "Model" in cname:
                    subtype = "Entity"
                
                if nid not in entities:
                    entities[nid] = {
                        "id": nid,
                        "label": cname,
                        "type": subtype,
                        "canonical_name": cname,
                        "aliases": [cname],
                        "provenance": []
                    }
                entities[nid]["provenance"].append({
                    "source_type": source_type,
                    "source_path": s3_path,
                    "source_chunk_id": cid,
                    "source_text": evidence_text[:500],
                    "confidence": 0.95
                })
                classes[cname] = nid
            
            # Interfaces
            for m in RE_JAVA_INTERFACE.finditer(text):
                iname = m.group(1)
                nid = node_id(f"Interface:{iname}")
                evidence_text = ctx(text, m.start(), 300)
                cid = find_chunk(m.start())
                
                if nid not in entities:
                    entities[nid] = {
                        "id": nid,
                        "label": iname,
                        "type": "Interface",
                        "canonical_name": iname,
                        "aliases": [iname],
                        "provenance": []
                    }
                entities[nid]["provenance"].append({
                    "source_type": source_type,
                    "source_path": s3_path,
                    "source_chunk_id": cid,
                    "source_text": evidence_text[:500],
                    "confidence": 0.93
                })
                interfaces[iname] = nid
            
            # Autowired dependencies -> relations
            for m in RE_AUTOWIRED.finditer(text):
                dep_name = m.group(1)
                evidence_text = ctx(text, m.start(), 200)
                cid = find_chunk(m.start())
                
                # Find the class that owns this autowired
                owner_class = None
                for cm in RE_JAVA_CLASS.finditer(text):
                    if cm.start() < m.start():
                        owner_class = cm.group(1)
                
                if owner_class:
                    src_id = classes.get(owner_class)
                    tgt_id = classes.get(dep_name) or interfaces.get(dep_name)
                    if src_id and tgt_id:
                        relations.append({
                            "from": src_id,
                            "to": tgt_id,
                            "type": "DEPENDS_ON",
                            "properties": {
                                "trigger_condition": "Spring DI injection",
                                "frequency": "runtime",
                                "protocol": "Spring_DI"
                            },
                            "provenance": [{
                                "source_type": source_type,
                                "source_path": s3_path,
                                "source_chunk_id": cid,
                                "source_text": evidence_text[:500],
                                "confidence": 0.92
                            }]
                        })
        
        # === Extract System references ===
        for m in RE_SYSTEM_REF.finditer(text):
            sys_name = m.group(1).upper()
            if sys_name in ('IMAPS', 'SUN', 'HDS', 'HULFT', 'ORACLE', 'MURATA', 'ERP', 'MDW'):
                nid = node_id(f"System:{sys_name}")
                evidence_text = ctx(text, m.start(), 200)
                cid = find_chunk(m.start())
                
                if nid not in entities:
                    entities[nid] = {
                        "id": nid,
                        "label": sys_name,
                        "type": "System",
                        "canonical_name": sys_name,
                        "aliases": [sys_name],
                        "provenance": []
                    }
                entities[nid]["provenance"].append({
                    "source_type": source_type,
                    "source_path": s3_path,
                    "source_chunk_id": cid,
                    "source_text": evidence_text[:500],
                    "confidence": 0.85
                })
                systems.add(sys_name)
        
        # === Extract INSERT...SELECT (ETL flows) ===
        if entry["extension"] in ('.sql', '.txt') and len(text) < 500000:
            for m in RE_INSERT_FROM.finditer(text[:100000]):
                target_table = m.group(1).upper()
                source_table = m.group(2).upper()
                evidence_text = ctx(text, m.start(), 400)
                cid = find_chunk(m.start())
                
                src_nid = tables.get(source_table, node_id(f"Table:{source_table}"))
                tgt_nid = tables.get(target_table, node_id(f"Table:{target_table}"))
                
                # Ensure both entities exist
                for tbl, nid in [(source_table, src_nid), (target_table, tgt_nid)]:
                    if nid not in entities:
                        entities[nid] = {
                            "id": nid,
                            "label": tbl,
                            "type": "Table",
                            "canonical_name": tbl,
                            "aliases": [tbl],
                            "provenance": [{
                                "source_type": source_type,
                                "source_path": s3_path,
                                "source_chunk_id": cid,
                                "source_text": evidence_text[:500],
                                "confidence": 0.90
                            }]
                        }
                
                relations.append({
                    "from": src_nid,
                    "to": tgt_nid,
                    "type": "FLOWS_TO",
                    "properties": {
                        "trigger_condition": "ETL INSERT...SELECT",
                        "frequency": "batch",
                        "protocol": "DB_Link"
                    },
                    "provenance": [{
                        "source_type": source_type,
                        "source_path": s3_path,
                        "source_chunk_id": cid,
                        "source_text": evidence_text[:500],
                        "confidence": 0.93
                    }]
                })
        
        # === Business Rules from column comments (STATUS patterns) ===
        for m in RE_COMMENT_COL.finditer(text):
            table = m.group(1).upper()
            col = m.group(2).upper()
            comment = m.group(3)
            evidence_text = ctx(text, m.start(), 300)
            cid = find_chunk(m.start())
            
            # Check for status/business rule
            if any(kw in comment for kw in ('状態', 'ステータス', 'STATUS', 'FLG', 'フラグ', '区分')) or re.search(r'\d+[：:]', comment):
                rule_id = node_id(f"Rule:{table}.{col}")
                if rule_id not in entities:
                    entities[rule_id] = {
                        "id": rule_id,
                        "label": f"{table}.{col}",
                        "type": "BusinessRule",
                        "canonical_name": f"{table}.{col}_rule",
                        "description": comment,
                        "aliases": [f"{table}.{col}"],
                        "provenance": [{
                            "source_type": source_type,
                            "source_path": s3_path,
                            "source_chunk_id": cid,
                            "source_text": evidence_text[:500],
                            "confidence": 0.90
                        }]
                    }
                
                # Link rule to table
                tbl_nid = tables.get(table)
                if tbl_nid:
                    relations.append({
                        "from": rule_id,
                        "to": tbl_nid,
                        "type": "GOVERNS",
                        "properties": {
                            "trigger_condition": f"Column {col} state change",
                            "frequency": "per_transaction",
                            "protocol": "DB_Constraint"
                        },
                        "provenance": [{
                            "source_type": source_type,
                            "source_path": s3_path,
                            "source_chunk_id": cid,
                            "source_text": evidence_text[:500],
                            "confidence": 0.88
                        }]
                    })
    
    # === Cross-document relationships (MVC pattern) ===
    print("  Building cross-document relations (MVC, Struts, config)...")
    
    # Action -> ServiceImpl -> DaoImpl pattern
    suffixes = {
        "Action": "Controller",
        "ServiceImpl": "Service", 
        "DaoImpl": "DataAccess",
        "Service": "Interface",
        "DaoI": "Interface"
    }
    
    # Group by base name
    base_groups = defaultdict(dict)
    for cname, nid in classes.items():
        for sfx in ["Action", "ServiceImpl", "DaoImpl"]:
            if cname.endswith(sfx):
                base = cname[:-len(sfx)]
                base_groups[base][sfx] = (cname, nid)
    for iname, nid in interfaces.items():
        for sfx in ["Service", "DaoI"]:
            if iname.endswith(sfx):
                base = iname[:-len(sfx)]
                base_groups[base][sfx] = (iname, nid)
    
    # Create MVC flow relations
    mvc_rules = [
        ("Action", "DELEGATES_TO", "ServiceImpl"),
        ("ServiceImpl", "CALLS_DAO", "DaoImpl"),
    ]
    
    for base, members in base_groups.items():
        for sfx_a, pred, sfx_b in mvc_rules:
            a = members.get(sfx_a)
            b = members.get(sfx_b)
            if a and b:
                relations.append({
                    "from": a[1],
                    "to": b[1],
                    "type": pred,
                    "properties": {
                        "trigger_condition": f"HTTP request to {base} endpoint",
                        "frequency": "per_request",
                        "protocol": "HTTP" if pred == "DELEGATES_TO" else "Method_Call"
                    },
                    "provenance": [{
                        "source_type": "code",
                        "source_path": f"s3://{BUCKET}/{PREFIX}代码_muratapr/",
                        "source_chunk_id": "cross_doc_inference",
                        "source_text": f"MVC pattern: {a[0]} {pred} {b[0]} based on naming convention",
                        "confidence": 0.88
                    }]
                })
    
    # === System-to-System relations from HDS SQL scripts ===
    # HDS processes data from iMaps to Murata MDW
    hds_sys = node_id("System:HDS")
    imaps_sys = node_id("System:IMAPS")
    mdw_sys = node_id("System:MDW")
    murata_sys = node_id("System:MURATA")
    
    # Ensure systems exist
    for sys_name, sys_id in [("HDS", hds_sys), ("IMAPS", imaps_sys), ("MDW", mdw_sys), ("MURATA", murata_sys)]:
        if sys_id not in entities:
            entities[sys_id] = {
                "id": sys_id,
                "label": sys_name,
                "type": "System",
                "canonical_name": sys_name,
                "aliases": [sys_name],
                "provenance": [{
                    "source_type": "design_doc",
                    "source_path": f"s3://{BUCKET}/{PREFIX}",
                    "source_chunk_id": "system_inference",
                    "source_text": f"System {sys_name} inferred from project structure and file naming",
                    "confidence": 0.90
                }]
            }
    
    # iMaps -> HDS -> MDW flow
    relations.append({
        "from": imaps_sys,
        "to": hds_sys,
        "type": "FLOWS_TO",
        "properties": {
            "trigger_condition": "Data sync from iMaps to HDS middleware",
            "frequency": "batch_daily",
            "protocol": "HULFT"
        },
        "provenance": [{
            "source_type": "sql",
            "source_path": f"s3://{BUCKET}/{PREFIX}HDS之SQL脚本/",
            "source_chunk_id": "system_flow_inference",
            "source_text": "HDS SQL scripts process data from iMaps (insert_journal_base references iMaps data fields)",
            "confidence": 0.85
        }]
    })
    
    relations.append({
        "from": hds_sys,
        "to": mdw_sys,
        "type": "FLOWS_TO",
        "properties": {
            "trigger_condition": "ETL processing complete",
            "frequency": "batch_daily",
            "protocol": "DB_Link"
        },
        "provenance": [{
            "source_type": "sql",
            "source_path": f"s3://{BUCKET}/{PREFIX}HDS之SQL脚本/",
            "source_chunk_id": "system_flow_inference",
            "source_text": "HDS scripts transform and load data into MDW database tables (journal_base, payment_req)",
            "confidence": 0.85
        }]
    })
    
    # === Document nodes ===
    for entry in manifest:
        if entry["is_empty"]:
            continue
        doc_id = node_id(f"Document:{entry['relative_path']}")
        entities[doc_id] = {
            "id": doc_id,
            "label": entry["file_name"],
            "type": "Document",
            "canonical_name": entry["file_name"],
            "aliases": [entry["file_name"]],
            "category": classify_doc_type(entry),
            "s3_path": entry["s3_path"],
            "relative_path": entry["relative_path"],
            "provenance": [{
                "source_type": classify_doc_type(entry),
                "source_path": entry["s3_path"],
                "source_chunk_id": chunks_by_file.get(entry["relative_path"], [{}])[0].get("chunk_id", ""),
                "source_text": f"File: {entry['file_name']} ({entry['char_count']} chars, {entry['extractor']} extractor)",
                "confidence": 1.0
            }]
        }
    
    # === Deduplicate relations ===
    seen_rels = set()
    unique_relations = []
    for r in relations:
        key = (r["from"], r["type"], r["to"])
        if key not in seen_rels:
            seen_rels.add(key)
            unique_relations.append(r)
        else:
            # Merge provenance
            for existing in unique_relations:
                if (existing["from"], existing["type"], existing["to"]) == key:
                    existing["provenance"].extend(r["provenance"])
                    # Multi-source increases confidence
                    max_conf = max(p["confidence"] for p in existing["provenance"])
                    for p in existing["provenance"]:
                        p["confidence"] = min(1.0, max_conf + 0.02)
                    break
    
    print(f"  Entities extracted: {len(entities)}")
    print(f"  Relations extracted: {len(unique_relations)}")
    print(f"  Systems found: {systems}")
    
    # Type breakdown
    type_counts = defaultdict(int)
    for e in entities.values():
        type_counts[e["type"]] += 1
    print("  Entity types:")
    for t, c in sorted(type_counts.items(), key=lambda x: -x[1]):
        print(f"    {t}: {c}")
    
    rel_type_counts = defaultdict(int)
    for r in unique_relations:
        rel_type_counts[r["type"]] += 1
    print("  Relation types:")
    for t, c in sorted(rel_type_counts.items(), key=lambda x: -x[1]):
        print(f"    {t}: {c}")
    
    return entities, unique_relations

# === Main ===
if __name__ == "__main__":
    start = time.time()
    
    # Step 1: Download
    files = download_s3_files()
    
    # Step 2: Extract text
    manifest = extract_texts(files)
    
    # Step 3: Chunk
    chunks = chunk_documents(manifest)
    
    # Step 4: Extract entities and relations
    entities, relations = extract_entities_and_relations(chunks, manifest)
    
    # Save intermediate results
    entities_path = OUTPUT_DIR / "entities_raw.json"
    relations_path = OUTPUT_DIR / "relations_raw.json"
    
    with open(entities_path, 'w', encoding='utf-8') as fp:
        json.dump(list(entities.values()), fp, ensure_ascii=False, indent=2)
    
    with open(relations_path, 'w', encoding='utf-8') as fp:
        json.dump(relations, fp, ensure_ascii=False, indent=2)
    
    elapsed = time.time() - start
    print(f"\n{'='*60}")
    print(f"Phase 1 Complete in {elapsed:.1f}s")
    print(f"  Files downloaded: {len(files)}")
    print(f"  Files extracted: {len(manifest)}")
    print(f"  Chunks generated: {len(chunks)}")
    print(f"  Entities: {len(entities)}")
    print(f"  Relations: {len(relations)}")
    print(f"  Output: {OUTPUT_DIR}")
    print(f"{'='*60}")
