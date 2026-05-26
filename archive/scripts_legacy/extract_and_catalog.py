#!/usr/bin/env python3
"""
Murata document extractor & metadata cataloger
================================================
Walks ~/hermes_graph_project/data/Murata/, extracts text from every file,
writes a per-file .txt sidecar under extracted/, and produces a master
manifest at ~/hermes_graph_project/data/manifest.json.

Supported formats
-----------------
  .txt .sql .csv .md .java .xml .jsp .properties .iml .css  -> plain read
  .pdf                                                        -> pdfminer.six
  .docx                                                       -> python-docx
  .xlsx .xls                                                  -> openpyxl / xlrd
  .pptx                                                       -> python-pptx
  everything else                                             -> try utf-8, skip binary
"""

import json
import os
import re
import sys
import traceback
from datetime import datetime, timezone
from pathlib import Path

# ── plain-text extensions ─────────────────────────────────────────────────────
PLAIN_TEXT_SUFFIXES = {
    ".txt", ".sql", ".csv", ".md", ".java", ".xml", ".jsp",
    ".properties", ".iml", ".css", ".MF", ".tern-project",
    ".classpath", ".project", ".pom", "",
}

ROOT      = Path.home() / "hermes_graph_project" / "data" / "Murata"
OUT_DIR   = Path.home() / "hermes_graph_project" / "data" / "extracted"
MANIFEST  = Path.home() / "hermes_graph_project" / "data" / "manifest.json"
S3_PREFIX = "s3://s3-hulftchina-rd/Murata"

# skip lock/temp files and pure binary assets unlikely to carry text
SKIP_PATTERNS = re.compile(
    r"(~\$|\.DS_Store|Thumbs\.db|\.jar$|\.gif$|\.png$|\.jpg$|\.db$"
    r"|pixel_0\.gif|NDH6SA)", re.IGNORECASE
)

OUT_DIR.mkdir(parents=True, exist_ok=True)


# ─────────────────────────────────────────────────────────────────────────────
# extractor helpers
# ─────────────────────────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> str:
    from pdfminer.high_level import extract_text
    try:
        text = extract_text(str(path))
        return text.strip() if text else ""
    except Exception as e:
        return f"[PDF extraction error: {e}]"


def extract_docx(path: Path) -> str:
    from docx import Document
    try:
        doc = Document(str(path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for shname in wb.sheetnames:
            ws = wb[shname]
            parts.append(f"=== Sheet: {shname} ===")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c not in (None, ""))
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[XLSX extraction error: {e}]"


def extract_xls(path: Path) -> str:
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        parts = []
        for sh in wb.sheets():
            parts.append(f"=== Sheet: {sh.name} ===")
            for rx in range(sh.nrows):
                row_text = " | ".join(str(sh.cell_value(rx, cx))
                                      for cx in range(sh.ncols)
                                      if str(sh.cell_value(rx, cx)).strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[XLS extraction error: {e}]"


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts)
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def extract_plain(path: Path) -> str:
    for enc in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            return path.read_text(encoding=enc).strip()
        except (UnicodeDecodeError, LookupError):
            continue
    return "[Could not decode as text — binary file skipped]"


def extract_text(path: Path) -> tuple[str, str]:
    """Return (extracted_text, method_used)."""
    suf = path.suffix.lower()
    if suf == ".pdf":
        return extract_pdf(path), "pdfminer"
    if suf == ".docx":
        return extract_docx(path), "python-docx"
    if suf == ".xlsx":
        return extract_xlsx(path), "openpyxl"
    if suf == ".xls":
        return extract_xls(path), "xlrd"
    if suf == ".pptx":
        return extract_pptx(path), "python-pptx"
    # plain text family
    if suf in PLAIN_TEXT_SUFFIXES or suf in {".java", ".xml", ".jsp", ".css"}:
        return extract_plain(path), "plain-text"
    # fallback — try plain text anyway
    text = extract_plain(path)
    method = "plain-text-fallback"
    return text, method


# ─────────────────────────────────────────────────────────────────────────────
# main
# ─────────────────────────────────────────────────────────────────────────────

def s3_key(local_path: Path) -> str:
    rel = local_path.relative_to(Path.home() / "hermes_graph_project" / "data")
    return f"s3://s3-hulftchina-rd/{rel.as_posix()}"


def build_manifest():
    records = []
    skipped = []
    errors  = []

    all_files = sorted(ROOT.rglob("*"))
    processable = [f for f in all_files if f.is_file() and not SKIP_PATTERNS.search(f.name)]

    print(f"Total files on disk  : {sum(1 for f in all_files if f.is_file())}")
    print(f"Files to process     : {len(processable)}")
    print(f"Files skipped (binary/temp): {sum(1 for f in all_files if f.is_file()) - len(processable)}")
    print()

    for fpath in processable:
        stat = fpath.stat()
        rel  = fpath.relative_to(ROOT)

        # sidecar output path
        out_path = OUT_DIR / (str(rel).replace("/", "__").replace("\\", "__") + ".txt")

        print(f"  [{fpath.suffix or 'no-ext':8s}] {rel.name[:60]}", end=" ... ", flush=True)

        try:
            text, method = extract_text(fpath)
            char_count = len(text)

            # write sidecar
            out_path.write_text(text, encoding="utf-8")

            record = {
                "file_name":    fpath.name,
                "relative_path": str(rel),
                "s3_path":      s3_key(fpath),
                "local_path":   str(fpath),
                "extracted_to": str(out_path),
                "size_bytes":   stat.st_size,
                "size_kb":      round(stat.st_size / 1024, 1),
                "modified_utc": datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc)
                                        .strftime("%Y-%m-%dT%H:%M:%SZ"),
                "extension":    fpath.suffix.lower() or "(none)",
                "extractor":    method,
                "char_count":   char_count,
                "is_empty":     char_count == 0,
            }
            records.append(record)
            print(f"OK  ({char_count:,} chars, {method})")

        except Exception:
            tb = traceback.format_exc()
            errors.append({"file": str(fpath), "error": tb})
            print(f"ERROR")

    # ── write manifest ─────────────────────────────────────────────────────────
    manifest = {
        "generated_at": datetime.now(tz=timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
        "s3_source":    "s3://s3-hulftchina-rd/Murata/",
        "local_root":   str(ROOT),
        "total_files":  len(records),
        "errors":       len(errors),
        "files":        records,
        "extraction_errors": errors,
    }
    MANIFEST.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    # ── summary ─────────────────────────────────────────────────────────────────
    print()
    print("=" * 60)
    print(" EXTRACTION COMPLETE")
    print("=" * 60)

    by_ext: dict[str, int] = {}
    for r in records:
        by_ext[r["extension"]] = by_ext.get(r["extension"], 0) + 1

    print(f"  Files processed     : {len(records)}")
    print(f"  Extraction errors   : {len(errors)}")
    print(f"  Empty extractions   : {sum(1 for r in records if r['is_empty'])}")
    print(f"  Manifest written to : {MANIFEST}")
    print(f"  Text sidecars dir   : {OUT_DIR}")
    print()
    print("  Breakdown by extension:")
    for ext, cnt in sorted(by_ext.items(), key=lambda x: -x[1]):
        print(f"    {ext:20s}  {cnt:3d} file(s)")
    print("=" * 60)

    return records, errors


if __name__ == "__main__":
    build_manifest()
