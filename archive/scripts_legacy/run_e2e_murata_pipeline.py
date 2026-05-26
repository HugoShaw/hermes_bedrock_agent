#!/usr/bin/env python3
"""Full E2E Murata Pipeline — scan → parse → chunk → embed → graph → load → retrieval → viz.

Usage:
    python scripts/run_e2e_murata_pipeline.py --stage all --confirm-live-write
    python scripts/run_e2e_murata_pipeline.py --stage scan
    python scripts/run_e2e_murata_pipeline.py --stage parse --enable-vlm
    python scripts/run_e2e_murata_pipeline.py --stage embedding --resume
    python scripts/run_e2e_murata_pipeline.py --stage graph --resume --skip-existing
    python scripts/run_e2e_murata_pipeline.py --stage load --live-neptune --confirm-live-write
    python scripts/run_e2e_murata_pipeline.py --stage retrieval
    python scripts/run_e2e_murata_pipeline.py --stage visualization

Supports:
- Full live run against real S3, Bedrock, LanceDB, Neptune
- Checkpoint/resume: re-run picks up where it left off
- Skip-existing: don't re-process documents already in JSONL artifacts
- Failure tracking: save per-file failures, continue on error
- Fail-fast: halt on first error if --fail-fast is set
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import sys
import time
import traceback
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional

# Ensure project is importable
sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "src"))

from hermes_bedrock_agent.configs.logging import get_logger

logger = get_logger("e2e_pipeline")


# ===========================================================================
# Configuration
# ===========================================================================

DEFAULT_RUN_ID = "murata_full_vlm_live_001"
DEFAULT_S3_BUCKET = "s3-hulftchina-rd"
DEFAULT_S3_PREFIX = "Murata/"
DEFAULT_S3_URI = f"s3://{DEFAULT_S3_BUCKET}/{DEFAULT_S3_PREFIX}"

DEFAULT_ARTIFACT_BASE = Path.home() / "projects/data/enterprise_graphrag/runs"
DEFAULT_LANCEDB_PATH = Path.home() / "projects/data/vector_store/lancedb"
DEFAULT_NEPTUNE_ENDPOINT = "g-nbuyck5yl8.ap-northeast-1.neptune-graph.amazonaws.com"

# File extensions to process (skip binary assets, .DS_Store, etc.)
PROCESSABLE_EXTENSIONS = {
    ".md", ".txt", ".sql", ".java", ".xml", ".properties",
    ".py", ".json", ".yaml", ".yml", ".csv", ".iml",
    ".jsp", ".css", ".html", ".htm",
    ".docx", ".pptx", ".xlsx", ".xls",
    ".pdf", ".png", ".jpg", ".jpeg", ".gif",
}

# Extensions to skip entirely
SKIP_EXTENSIONS = {
    ".DS_Store", ".jar", ".class", ".db", ".Thumbs.db",
}

STAGE_ORDER = ["scan", "parse", "chunk", "embedding", "graph", "load", "enrichment", "retrieval", "visualization"]


# ===========================================================================
# Pipeline State
# ===========================================================================


class PipelineState:
    """Manages pipeline run state, artifacts directory, and checkpoint files."""

    def __init__(self, run_id: str, artifact_base: Path):
        self.run_id = run_id
        self.artifact_dir = artifact_base / run_id / "artifacts"
        self.artifact_dir.mkdir(parents=True, exist_ok=True)
        self.failures: list[dict] = []
        self.start_time = time.time()

    def artifact_path(self, filename: str) -> Path:
        return self.artifact_dir / filename

    def load_jsonl(self, filename: str) -> list[dict]:
        """Load existing JSONL artifact (for resume)."""
        path = self.artifact_path(filename)
        if not path.exists():
            return []
        records = []
        with open(path, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if line:
                    records.append(json.loads(line))
        return records

    def save_jsonl(self, filename: str, records: list[dict | Any]) -> int:
        """Save records as JSONL. Returns count written."""
        path = self.artifact_path(filename)
        with open(path, "w", encoding="utf-8") as f:
            for r in records:
                if hasattr(r, "model_dump"):
                    data = r.model_dump(mode="json")
                elif hasattr(r, "dict"):
                    data = r.dict()
                elif isinstance(r, dict):
                    data = r
                else:
                    data = {"value": str(r)}
                f.write(json.dumps(data, ensure_ascii=False, default=str) + "\n")
        logger.info(f"Saved {len(records)} records to {path}")
        return len(records)

    def append_jsonl(self, filename: str, record: dict | Any) -> None:
        """Append a single record to a JSONL file."""
        path = self.artifact_path(filename)
        with open(path, "a", encoding="utf-8") as f:
            if hasattr(record, "model_dump"):
                data = record.model_dump(mode="json")
            elif isinstance(record, dict):
                data = record
            else:
                data = {"value": str(record)}
            f.write(json.dumps(data, ensure_ascii=False, default=str) + "\n")

    def save_json(self, filename: str, data: dict | list) -> None:
        """Save a single JSON report file."""
        path = self.artifact_path(filename)
        with open(path, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, ensure_ascii=False, default=str)
        logger.info(f"Saved report: {path}")

    def save_text(self, filename: str, text: str) -> None:
        """Save a text/markdown/cypher artifact."""
        path = self.artifact_path(filename)
        with open(path, "w", encoding="utf-8") as f:
            f.write(text)
        logger.info(f"Saved text: {path}")

    def record_failure(self, stage: str, item_id: str, error: str) -> None:
        """Record a processing failure."""
        self.failures.append({
            "stage": stage,
            "item_id": item_id,
            "error": error,
            "timestamp": datetime.now(timezone.utc).isoformat(),
        })

    def get_existing_ids(self, filename: str, id_field: str) -> set[str]:
        """Load existing IDs from a JSONL file (for skip-existing)."""
        records = self.load_jsonl(filename)
        return {r.get(id_field, "") for r in records if r.get(id_field)}


# ===========================================================================
# Stage 1: SCAN
# ===========================================================================


def stage_scan(state: PipelineState, args) -> list[dict]:
    """Scan S3 bucket and create document inventory."""
    logger.info("=" * 60)
    logger.info("STAGE 1: SCAN — S3 document discovery")
    logger.info("=" * 60)

    from hermes_bedrock_agent.clients.s3_client import S3Client, S3Object
    from hermes_bedrock_agent.schemas.document import SourceDocument, SourceType
    from hermes_bedrock_agent.utils.hashing import content_hash as compute_hash

    # Parse S3 URI from args
    s3_uri = getattr(args, "s3_uri", DEFAULT_S3_URI)
    if s3_uri.startswith("s3://"):
        parts = s3_uri[5:].split("/", 1)
        bucket = parts[0]
        prefix = parts[1] if len(parts) > 1 else ""
    else:
        bucket = DEFAULT_S3_BUCKET
        prefix = DEFAULT_S3_PREFIX

    s3 = S3Client(bucket=bucket)
    logger.info(f"Scanning s3://{bucket}/{prefix} ...")

    # List all objects
    objects = s3.list_objects(prefix=prefix)
    logger.info(f"Found {len(objects)} total objects in S3")

    # Filter processable files
    documents = []
    skipped = []
    for obj in objects:
        ext = Path(obj.key).suffix.lower()
        filename = Path(obj.key).name

        # Skip binary/system files
        if filename in (".DS_Store", "Thumbs.db") or ext in SKIP_EXTENSIONS:
            skipped.append({"key": obj.key, "reason": "system_file"})
            continue
        if obj.size == 0:
            skipped.append({"key": obj.key, "reason": "empty"})
            continue
        if filename.startswith("~$"):
            skipped.append({"key": obj.key, "reason": "temp_file"})
            continue

        # Determine source type
        source_type = _infer_source_type(ext, obj.key)

        # Build document record
        doc_id = hashlib.sha256(obj.uri.encode()).hexdigest()[:16]
        doc = {
            "document_id": doc_id,
            "source_uri": obj.uri,
            "source_type": source_type,
            "filename": filename,
            "file_size": obj.size,
            "s3_key": obj.key,
            "etag": obj.etag,
            "last_modified": obj.last_modified.isoformat() if obj.last_modified else "",
            "run_id": state.run_id,
            "dataset": getattr(args, "neptune_dataset", "murata"),
            "source_prefix": s3_uri,
        }
        documents.append(doc)

    # Save artifacts
    state.save_jsonl("documents.jsonl", documents)

    # Build inventory report
    type_counts: dict[str, int] = {}
    for doc in documents:
        t = doc["source_type"]
        type_counts[t] = type_counts.get(t, 0) + 1

    inventory = {
        "run_id": state.run_id,
        "s3_uri": s3_uri,
        "scan_time": datetime.now(timezone.utc).isoformat(),
        "total_objects": len(objects),
        "processable_documents": len(documents),
        "skipped_count": len(skipped),
        "type_distribution": type_counts,
        "skipped_details": skipped[:20],  # First 20 only
        "total_bytes": sum(d["file_size"] for d in documents),
    }
    state.save_json("file_inventory_report.json", inventory)

    logger.info(f"Scan complete: {len(documents)} processable documents")
    logger.info(f"Type distribution: {type_counts}")
    return documents


def _infer_source_type(ext: str, key: str) -> str:
    """Infer document source type from extension and path."""
    ext_map = {
        ".md": "markdown", ".txt": "text",
        ".sql": "sql", ".SQL": "sql",
        ".java": "code", ".py": "code", ".xml": "config",
        ".properties": "config", ".yml": "config", ".yaml": "config",
        ".json": "config", ".iml": "config",
        ".jsp": "code", ".css": "code", ".html": "code", ".htm": "code",
        ".csv": "text", ".CSV": "text",
        ".docx": "docx", ".pptx": "pptx",
        ".xlsx": "spreadsheet", ".xls": "spreadsheet",
        ".pdf": "pdf",
        ".png": "image", ".jpg": "image", ".jpeg": "image", ".gif": "image",
        ".classpath": "config", ".project": "config",
    }
    return ext_map.get(ext, "unknown")


# ===========================================================================
# Stage 2: PARSE
# ===========================================================================


def stage_parse(state: PipelineState, args) -> tuple[list[dict], list[dict]]:
    """Parse documents: text extraction + VLM for images/PDF pages."""
    logger.info("=" * 60)
    logger.info("STAGE 2: PARSE — Document parsing + VLM")
    logger.info("=" * 60)

    from hermes_bedrock_agent.clients.s3_client import S3Client
    from hermes_bedrock_agent.schemas.document import NormalizedDocument, SourceType

    s3 = S3Client(bucket=DEFAULT_S3_BUCKET)

    # Load scan results
    documents = state.load_jsonl("documents.jsonl")
    if not documents:
        raise RuntimeError("No documents.jsonl found. Run scan stage first.")

    # Check for resume
    existing_doc_ids = set()
    if args.skip_existing:
        existing_doc_ids = state.get_existing_ids("normalized_documents.jsonl", "document_id")
        logger.info(f"Skip-existing: {len(existing_doc_ids)} already parsed")

    normalized_docs = []
    visual_blocks = []
    parser_failures = []
    bedrock_client = None

    if args.enable_vlm:
        import boto3
        bedrock_client = boto3.client("bedrock-runtime", region_name="ap-northeast-1")

    for i, doc in enumerate(documents):
        doc_id = doc["document_id"]

        if doc_id in existing_doc_ids:
            continue

        try:
            logger.info(f"[{i+1}/{len(documents)}] Parsing: {doc['filename']} ({doc['source_type']})")

            # Download from S3
            raw_bytes = s3.download_bytes(doc["s3_key"])

            # Parse based on type
            source_type = doc["source_type"]
            sections = []
            doc_visual_blocks = []

            if source_type in ("markdown", "text", "code", "config", "sql"):
                sections = _parse_text(raw_bytes, doc)
            elif source_type == "docx":
                sections = _parse_docx(raw_bytes, doc)
            elif source_type == "pptx":
                sections = _parse_pptx(raw_bytes, doc)
            elif source_type == "spreadsheet":
                sections = _parse_spreadsheet(raw_bytes, doc)
            elif source_type == "pdf":
                sections, doc_visual_blocks = _parse_pdf(
                    raw_bytes, doc, bedrock_client if args.enable_vlm else None
                )
            elif source_type == "image":
                if args.enable_vlm and bedrock_client:
                    doc_visual_blocks = _parse_image_vlm(raw_bytes, doc, bedrock_client)
                    # Fallback: if VLM fails, still register image with placeholder
                    if not doc_visual_blocks:
                        sections = [{"title": "Image", "content": f"[Image file: {doc['filename']} — VLM analysis failed]", "page": ""}]
                else:
                    sections = [{"title": "Image", "content": f"[Image file: {doc['filename']}]", "page": ""}]
            else:
                sections = _parse_text(raw_bytes, doc)  # Fallback

            # Build NormalizedDocument
            norm_doc = {
                "document_id": doc_id,
                "source_uri": doc["source_uri"],
                "source_type": source_type,
                "filename": doc["filename"],
                "title": doc["filename"],
                "sections": sections,
                "total_pages": len({s.get("page") for s in sections if s.get("page")}) or None,
                "language": "ja",  # Murata is Japanese enterprise
                "visual_block_ids": [vb["visual_id"] for vb in doc_visual_blocks],
                "content_hash": hashlib.sha256(raw_bytes).hexdigest(),
                "run_id": state.run_id,
                "dataset": "murata",
                "source_prefix": DEFAULT_S3_URI,
                "parsed_at": datetime.now(timezone.utc).isoformat(),
            }
            normalized_docs.append(norm_doc)
            visual_blocks.extend(doc_visual_blocks)

            # Save incrementally (for resume)
            state.append_jsonl("normalized_documents.jsonl", norm_doc)
            for vb in doc_visual_blocks:
                state.append_jsonl("visual_blocks.jsonl", vb)

        except Exception as e:
            error_msg = f"{type(e).__name__}: {e}"
            logger.error(f"Failed to parse {doc['filename']}: {error_msg}")
            failure = {
                "document_id": doc_id,
                "filename": doc["filename"],
                "source_type": doc["source_type"],
                "error": error_msg,
                "traceback": traceback.format_exc(),
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }
            parser_failures.append(failure)
            state.record_failure("parse", doc_id, error_msg)

            if args.fail_fast:
                raise

    # Save failures
    if parser_failures:
        state.save_jsonl("parser_failed.jsonl", parser_failures)
    else:
        state.save_jsonl("parser_failed.jsonl", [])

    logger.info(f"Parse complete: {len(normalized_docs)} docs, "
                f"{len(visual_blocks)} visual blocks, {len(parser_failures)} failures")
    return normalized_docs, visual_blocks


def _parse_text(raw_bytes: bytes, doc: dict) -> list[dict]:
    """Parse text/code/markdown/sql files."""
    # Try multiple encodings
    text = None
    for encoding in ("utf-8", "shift_jis", "gbk", "latin-1"):
        try:
            text = raw_bytes.decode(encoding)
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if text is None:
        text = raw_bytes.decode("utf-8", errors="replace")

    # Split into sections by headers or logical boundaries
    sections = []
    if doc["source_type"] == "markdown":
        sections = _split_markdown_sections(text)
    elif doc["source_type"] == "sql":
        sections = _split_sql_sections(text)
    elif doc["source_type"] == "code":
        sections = _split_code_sections(text, doc["filename"])
    else:
        sections = [{"title": doc["filename"], "content": text, "page": ""}]
    return sections


def _split_markdown_sections(text: str) -> list[dict]:
    """Split markdown by headers."""
    import re
    sections = []
    current_title = "Introduction"
    current_content = []

    for line in text.split("\n"):
        header_match = re.match(r"^(#{1,4})\s+(.+)$", line)
        if header_match:
            if current_content:
                content = "\n".join(current_content).strip()
                if content:
                    sections.append({"title": current_title, "content": content, "page": ""})
            current_title = header_match.group(2)
            current_content = []
        else:
            current_content.append(line)

    if current_content:
        content = "\n".join(current_content).strip()
        if content:
            sections.append({"title": current_title, "content": content, "page": ""})

    return sections or [{"title": "Document", "content": text, "page": ""}]


def _split_sql_sections(text: str) -> list[dict]:
    """Split SQL by statement boundaries."""
    import re
    statements = re.split(r";\s*\n", text)
    sections = []
    for i, stmt in enumerate(statements):
        stmt = stmt.strip()
        if not stmt or len(stmt) < 5:
            continue
        # Try to identify the object name
        title_match = re.search(
            r"(?:CREATE|ALTER|INSERT INTO|UPDATE)\s+(?:TABLE|VIEW|INDEX|PROCEDURE)?\s*`?(\w+)`?",
            stmt, re.IGNORECASE
        )
        title = title_match.group(1) if title_match else f"Statement_{i+1}"
        sections.append({"title": title, "content": stmt, "page": ""})
    return sections or [{"title": "SQL", "content": text, "page": ""}]


def _split_code_sections(text: str, filename: str) -> list[dict]:
    """Split Java/code by class/method boundaries."""
    import re
    # For Java, split by class/method
    if filename.endswith(".java"):
        sections = []
        # Find class-level split
        class_match = re.search(r"(?:public|private)?\s*class\s+(\w+)", text)
        class_name = class_match.group(1) if class_match else filename

        # Split by methods
        method_pattern = r"(?:public|private|protected)\s+\w+[\w<>\[\],\s]*\s+(\w+)\s*\("
        methods = list(re.finditer(method_pattern, text))

        if methods:
            # Prologue (imports, class header)
            prologue = text[:methods[0].start()].strip()
            if prologue:
                sections.append({"title": f"{class_name} (imports/header)", "content": prologue, "page": ""})

            for j, m in enumerate(methods):
                start = m.start()
                end = methods[j+1].start() if j+1 < len(methods) else len(text)
                method_text = text[start:end].strip()
                sections.append({"title": f"{class_name}.{m.group(1)}()", "content": method_text, "page": ""})
        else:
            sections = [{"title": class_name, "content": text, "page": ""}]
        return sections
    else:
        return [{"title": filename, "content": text, "page": ""}]


def _parse_docx(raw_bytes: bytes, doc: dict) -> list[dict]:
    """Parse DOCX using python-docx."""
    try:
        import io
        from docx import Document as DocxDocument
        docx_doc = DocxDocument(io.BytesIO(raw_bytes))

        sections = []
        current_section = {"title": doc["filename"], "content": "", "page": ""}

        for para in docx_doc.paragraphs:
            if para.style.name.startswith("Heading"):
                if current_section["content"].strip():
                    sections.append(current_section)
                current_section = {"title": para.text or "Section", "content": "", "page": ""}
            else:
                current_section["content"] += para.text + "\n"

        if current_section["content"].strip():
            sections.append(current_section)

        return sections or [{"title": doc["filename"], "content": "[Empty DOCX]", "page": ""}]
    except ImportError:
        logger.warning("python-docx not installed, falling back to raw extract")
        return [{"title": doc["filename"], "content": f"[DOCX: {doc['filename']} - requires python-docx]", "page": ""}]
    except Exception as e:
        logger.warning(f"DOCX parse error: {e}")
        return [{"title": doc["filename"], "content": f"[DOCX parse failed: {e}]", "page": ""}]


def _parse_pptx(raw_bytes: bytes, doc: dict) -> list[dict]:
    """Parse PPTX using python-pptx."""
    try:
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw_bytes))

        sections = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                content = "\n".join(texts)
                sections.append({"title": f"Slide {slide_num}", "content": content, "page": str(slide_num)})

        return sections or [{"title": doc["filename"], "content": "[Empty PPTX]", "page": ""}]
    except ImportError:
        return [{"title": doc["filename"], "content": f"[PPTX: {doc['filename']} - requires python-pptx]", "page": ""}]
    except Exception as e:
        return [{"title": doc["filename"], "content": f"[PPTX parse failed: {e}]", "page": ""}]


def _parse_spreadsheet(raw_bytes: bytes, doc: dict) -> list[dict]:
    """Parse Excel spreadsheet."""
    try:
        import io
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)

        sections = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(max_row=200, values_only=True):
                row_text = " | ".join(str(c) if c is not None else "" for c in row)
                if row_text.strip(" |"):
                    rows.append(row_text)
            if rows:
                content = "\n".join(rows[:100])  # Limit rows
                sections.append({"title": f"Sheet: {sheet_name}", "content": content, "page": ""})
        wb.close()
        return sections or [{"title": doc["filename"], "content": "[Empty spreadsheet]", "page": ""}]
    except ImportError:
        return [{"title": doc["filename"], "content": f"[Spreadsheet: {doc['filename']} - requires openpyxl]", "page": ""}]
    except Exception as e:
        return [{"title": doc["filename"], "content": f"[Spreadsheet parse failed: {e}]", "page": ""}]


def _parse_pdf(raw_bytes: bytes, doc: dict, bedrock_client=None) -> tuple[list[dict], list[dict]]:
    """Parse PDF: text layer + optional VLM for page images."""
    sections = []
    visual_blocks = []

    try:
        import io
        import fitz  # PyMuPDF
        pdf_doc = fitz.open(stream=raw_bytes, filetype="pdf")

        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            text = page.get_text("text")
            if text.strip():
                sections.append({
                    "title": f"Page {page_num + 1}",
                    "content": text,
                    "page": str(page_num + 1),
                })

            # VLM: render page as image and analyze
            if bedrock_client and page_num < 20:  # Limit VLM to first 20 pages
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom
                    img_bytes = pix.tobytes("png")
                    vb = _vlm_analyze_image(
                        img_bytes, doc, bedrock_client,
                        context=f"PDF page {page_num + 1} of {doc['filename']}",
                        page=page_num + 1,
                    )
                    if vb:
                        visual_blocks.append(vb)
                except Exception as e:
                    logger.warning(f"VLM failed for page {page_num+1}: {e}")

        pdf_doc.close()
    except ImportError:
        sections = [{"title": doc["filename"], "content": "[PDF: requires PyMuPDF (fitz)]", "page": ""}]
    except Exception as e:
        sections = [{"title": doc["filename"], "content": f"[PDF parse failed: {e}]", "page": ""}]

    return sections, visual_blocks


def _parse_image_vlm(raw_bytes: bytes, doc: dict, bedrock_client) -> list[dict]:
    """Parse image via VLM (Bedrock Claude multimodal)."""
    vb = _vlm_analyze_image(
        raw_bytes, doc, bedrock_client,
        context=f"Image file: {doc['filename']}",
        page=None,
    )
    return [vb] if vb else []


def _vlm_analyze_image(
    image_bytes: bytes,
    doc: dict,
    bedrock_client,
    context: str = "",
    page: Optional[int] = None,
) -> Optional[dict]:
    """Call Bedrock Claude multimodal to analyze an image."""
    import base64

    # Determine media type
    filename = doc["filename"].lower()
    if filename.endswith(".png"):
        media_type = "image/png"
    elif filename.endswith((".jpg", ".jpeg")):
        media_type = "image/jpeg"
    elif filename.endswith(".gif"):
        media_type = "image/gif"
    else:
        media_type = "image/png"  # Default for PDF page renders

    img_b64 = base64.b64encode(image_bytes).decode("utf-8")

    # Size guard: Claude has 5MB image limit; skip very large images
    if len(image_bytes) > 4_500_000:
        logger.warning(f"Image too large for VLM ({len(image_bytes)} bytes), skipping")
        return None

    prompt = (
        f"Analyze this image from an enterprise system document. Context: {context}\n\n"
        "Please describe:\n"
        "1. What type of diagram/screenshot/document this is\n"
        "2. Key entities, systems, or components shown\n"
        "3. Relationships or data flows between entities\n"
        "4. Any text, labels, or annotations visible\n"
        "5. Business meaning in the context of enterprise AP/payment systems\n\n"
        "Output in structured format with sections. Be precise about entity names."
    )

    try:
        response = bedrock_client.converse(
            modelId="apac.anthropic.claude-sonnet-4-20250514-v1:0",
            messages=[{
                "role": "user",
                "content": [
                    {"image": {"format": media_type.split("/")[1], "source": {"bytes": image_bytes}}},
                    {"text": prompt},
                ],
            }],
            inferenceConfig={"maxTokens": 2048, "temperature": 0.0},
        )

        # Extract text
        output = response.get("output", {})
        message = output.get("message", {})
        content_blocks = message.get("content", [])
        description = ""
        for block in content_blocks:
            if "text" in block:
                description += block["text"]

        if not description.strip():
            return None

        visual_id = hashlib.sha256(
            f"{doc['document_id']}:page={page}:vlm".encode()
        ).hexdigest()[:16]

        return {
            "visual_id": visual_id,
            "document_id": doc["document_id"],
            "source_uri": doc.get("source_uri", ""),
            "page": page or 0,
            "visual_type": "page_screenshot",
            "visual_summary": description,
            "extracted_text": description[:500],
            "model_name": "apac.anthropic.claude-sonnet-4-20250514-v1:0",
            "confidence": 0.85,
            "run_id": doc.get("run_id", ""),
            "created_at": datetime.now(timezone.utc).isoformat(),
        }

    except Exception as e:
        logger.warning(f"VLM call failed: {e}")
        return None


# ===========================================================================
# Stage 3: CHUNK
# ===========================================================================


def stage_chunk(state: PipelineState, args) -> list[dict]:
    """Chunk normalized documents into DocumentChunks."""
    logger.info("=" * 60)
    logger.info("STAGE 3: CHUNK — Structure-aware chunking")
    logger.info("=" * 60)

    from hermes_bedrock_agent.chunking.chunker import ChunkerConfig, StructureAwareChunker
    from hermes_bedrock_agent.schemas.chunk import DocumentChunk
    from hermes_bedrock_agent.schemas.document import (
        NormalizedDocument,
        SourceType,
    )

    # Load normalized documents
    norm_docs = state.load_jsonl("normalized_documents.jsonl")
    visual_blocks_raw = state.load_jsonl("visual_blocks.jsonl")
    if not norm_docs:
        raise RuntimeError("No normalized_documents.jsonl found. Run parse stage first.")

    chunker = StructureAwareChunker(config=ChunkerConfig(
        chunk_size=1500,
        chunk_overlap=200,
        max_chunk_size=3000,
        code_chunk_size=2000,
        include_visual_chunks=True,
    ))

    all_chunks = []
    for doc_data in norm_docs:
        try:
            # Build NormalizedDocument from dict
            # sections are list[dict[str, str]] in NormalizedDocument schema
            sections = doc_data.get("sections", [])

            # Chunker uses document.content (single string), not sections directly
            # Concatenate section content into the content field
            content_parts = []
            for s in sections:
                if s.get("title"):
                    content_parts.append(f"## {s['title']}")
                if s.get("content"):
                    content_parts.append(s["content"])
            full_content = "\n\n".join(content_parts)

            source_type_str = doc_data.get("source_type", "unknown")
            try:
                source_type = SourceType(source_type_str)
            except ValueError:
                source_type = SourceType.UNKNOWN

            norm_doc = NormalizedDocument(
                document_id=doc_data["document_id"],
                source_uri=doc_data.get("source_uri", ""),
                source_type=source_type,
                title=doc_data.get("title", doc_data.get("filename", "")),
                content=full_content,
                sections=sections,
                page_count=doc_data.get("total_pages") or 0,
                language=doc_data.get("language", ""),
                visual_block_ids=doc_data.get("visual_block_ids", []),
                content_hash=doc_data.get("content_hash", ""),
            )

            # Get visual blocks for this doc
            doc_vbs = [
                vb for vb in visual_blocks_raw
                if vb.get("document_id") == doc_data["document_id"]
            ]

            # Import VisualBlock schema
            from hermes_bedrock_agent.schemas.visual import VisualBlock as VBModel
            vb_models = []
            for vb in doc_vbs:
                try:
                    vb_models.append(VBModel(
                        visual_id=vb.get("visual_id", ""),
                        document_id=vb.get("document_id", ""),
                        source_uri=vb.get("source_uri", ""),
                        page=vb.get("page", 0),
                        visual_type=vb.get("visual_type", "unknown"),
                        visual_summary=vb.get("visual_summary", vb.get("description", "")),
                        extracted_text=vb.get("extracted_text", ""),
                    ))
                except Exception:
                    pass

            chunks = chunker.chunk_document(norm_doc, visual_blocks=vb_models)
            all_chunks.extend(chunks)

        except Exception as e:
            logger.error(f"Chunking failed for {doc_data.get('filename', '?')}: {e}")
            state.record_failure("chunk", doc_data.get("document_id", ""), str(e))
            if args.fail_fast:
                raise

    # Save chunks
    chunk_dicts = []
    for c in all_chunks:
        d = c.model_dump(mode="json") if hasattr(c, "model_dump") else c.__dict__
        d["run_id"] = state.run_id
        d["dataset"] = "murata"
        d["source_prefix"] = DEFAULT_S3_URI
        chunk_dicts.append(d)

    state.save_jsonl("chunks.jsonl", chunk_dicts)

    # Stats
    stats = {
        "run_id": state.run_id,
        "total_documents": len(norm_docs),
        "total_chunks": len(all_chunks),
        "avg_chunks_per_doc": round(len(all_chunks) / max(len(norm_docs), 1), 2),
        "chunk_type_distribution": {},
        "timestamp": datetime.now(timezone.utc).isoformat(),
    }
    type_dist: dict[str, int] = {}
    for c in chunk_dicts:
        ct = c.get("chunk_type", "unknown")
        type_dist[ct] = type_dist.get(ct, 0) + 1
    stats["chunk_type_distribution"] = type_dist
    state.save_json("chunk_stats.json", stats)

    logger.info(f"Chunking complete: {len(all_chunks)} chunks from {len(norm_docs)} docs")
    return chunk_dicts


# ===========================================================================
# Stage 4: EMBEDDING
# ===========================================================================


def _matches_exclude_patterns(chunk: dict, patterns: list[str], docs_by_id: dict) -> bool:
    """Check if a chunk matches any of the exclude glob patterns.

    Matches against: filename, source_uri, document_id, section_title.
    Uses fnmatch for glob-style patterns (*, ?, []).
    """
    import fnmatch

    if not patterns:
        return False

    doc = docs_by_id.get(chunk.get("document_id", ""), {})
    targets = [
        doc.get("filename", ""),
        chunk.get("source_uri", "") or doc.get("source_uri", ""),
        chunk.get("document_id", ""),
        chunk.get("section_title", ""),
    ]

    for pattern in patterns:
        for target in targets:
            if target and fnmatch.fnmatch(target, pattern):
                return True
    return False


def _filter_chunks_for_stage(
    chunks: list[dict],
    docs_by_id: dict,
    stage_name: str,
    exclude_patterns: list[str],
    stage_patterns: list[str],
) -> tuple[list[dict], dict]:
    """Filter chunks based on exclude patterns for a given stage.

    Args:
        chunks: All chunks from chunks.jsonl
        docs_by_id: Document lookup dict {document_id: doc_dict}
        stage_name: 'embedding' or 'graph_extraction'
        exclude_patterns: Global --exclude-pattern list
        stage_patterns: Stage-specific --skip-*-pattern list

    Returns:
        (filtered_chunks, filter_summary_dict)
    """
    all_patterns = exclude_patterns + stage_patterns

    if not all_patterns:
        return chunks, {
            "stage": stage_name,
            "original_count": len(chunks),
            "excluded_count": 0,
            "included_count": len(chunks),
            "excluded_files": [],
            "patterns_applied": [],
        }

    included = []
    excluded = []
    excluded_files = set()

    for chunk in chunks:
        if _matches_exclude_patterns(chunk, all_patterns, docs_by_id):
            excluded.append(chunk)
            doc = docs_by_id.get(chunk.get("document_id", ""), {})
            excluded_files.add(doc.get("filename", chunk.get("document_id", "unknown")))
        else:
            included.append(chunk)

    summary = {
        "stage": stage_name,
        "original_count": len(chunks),
        "excluded_count": len(excluded),
        "included_count": len(included),
        "excluded_files": sorted(excluded_files),
        "patterns_applied": all_patterns,
    }

    logger.info(
        f"Filter ({stage_name}): {len(chunks)} total → "
        f"{len(included)} included, {len(excluded)} excluded "
        f"({len(excluded_files)} files matched patterns {all_patterns})"
    )

    return included, summary


def stage_embedding(state: PipelineState, args) -> list[dict]:
    """Generate embeddings and write to LanceDB."""
    logger.info("=" * 60)
    logger.info("STAGE 4: EMBEDDING — Bedrock Titan + LanceDB write")
    logger.info("=" * 60)

    from hermes_bedrock_agent.embedding.embedder import BedrockEmbedder, EmbedderConfig
    from hermes_bedrock_agent.schemas.chunk import ChunkEmbedding, DocumentChunk, ChunkType
    from hermes_bedrock_agent.vector_store import create_vector_store
    from hermes_bedrock_agent.vector_store.base import VectorStoreRecord

    # Load chunks
    chunk_dicts = state.load_jsonl("chunks.jsonl")
    if not chunk_dicts:
        raise RuntimeError("No chunks.jsonl found. Run chunk stage first.")

    # Apply exclude patterns for embedding
    docs_by_id = {}
    doc_dicts = state.load_jsonl("normalized_documents.jsonl")
    if doc_dicts:
        docs_by_id = {d["document_id"]: d for d in doc_dicts}

    exclude_patterns = getattr(args, "exclude_pattern", []) or []
    skip_emb_patterns = getattr(args, "skip_embedding_pattern", []) or []
    chunk_dicts, filter_summary = _filter_chunks_for_stage(
        chunk_dicts, docs_by_id, "embedding", exclude_patterns, skip_emb_patterns
    )
    state.save_jsonl("embedding_filter_summary.json", [filter_summary])

    # Skip-existing
    existing_ids = set()
    if args.skip_existing:
        existing_ids = state.get_existing_ids("embeddings.jsonl", "chunk_id")
        logger.info(f"Skip-existing: {len(existing_ids)} already embedded")

    chunks_to_embed = [c for c in chunk_dicts if c.get("chunk_id") not in existing_ids]
    logger.info(f"Embedding {len(chunks_to_embed)}/{len(chunk_dicts)} chunks")

    # Initialize embedder
    embedder = BedrockEmbedder(config=EmbedderConfig(
        model_id="amazon.titan-embed-text-v2:0",
        dimension=1024,
        batch_size=10,
        max_retries=3,
    ))

    # Embed in batches
    embeddings = []
    lancedb_path = str(getattr(args, "local_vector_store_path", None) or DEFAULT_LANCEDB_PATH)
    collection_name = getattr(args, "local_vector_collection", None) or f"murata_e2e_{state.run_id}"
    use_mock_embedding = getattr(args, "mock_embedding", False)

    for i in range(0, len(chunks_to_embed), 10):
        batch = chunks_to_embed[i:i+10]
        for chunk_dict in batch:
            try:
                # Build DocumentChunk for embedder
                chunk_type_str = chunk_dict.get("chunk_type", "text")
                try:
                    chunk_type = ChunkType(chunk_type_str)
                except (ValueError, KeyError):
                    chunk_type = ChunkType.TEXT

                doc_chunk = DocumentChunk(
                    chunk_id=chunk_dict["chunk_id"],
                    document_id=chunk_dict.get("document_id", ""),
                    content=chunk_dict.get("content", ""),
                    chunk_type=chunk_type,
                    chunk_index=chunk_dict.get("chunk_index", 0),
                    source_uri=chunk_dict.get("source_uri", ""),
                    source_type=chunk_dict.get("source_type", "unknown"),
                    section_title=chunk_dict.get("section_title", ""),
                    page=chunk_dict.get("page"),
                    content_hash=chunk_dict.get("content_hash", ""),
                    visual_block_ids=chunk_dict.get("visual_block_ids", []),
                    acl=chunk_dict.get("acl", []),
                )

                if use_mock_embedding:
                    import random
                    emb_dict = {
                        "chunk_id": doc_chunk.chunk_id,
                        "document_id": doc_chunk.document_id,
                        "content": doc_chunk.content[:200],
                        "embedding": [random.gauss(0, 0.1) for _ in range(1024)],
                        "embedding_model": "mock-random-1024d",
                        "embedding_dimension": 1024,
                        "run_id": state.run_id,
                    }
                else:
                    emb = embedder.embed_chunk(doc_chunk)
                    emb_dict = emb.model_dump(mode="json")
                    emb_dict["run_id"] = state.run_id

                embeddings.append(emb_dict)

                # Incremental save
                state.append_jsonl("embeddings.jsonl", emb_dict)

            except Exception as e:
                logger.error(f"Embedding failed for {chunk_dict.get('chunk_id', '?')}: {e}")
                state.record_failure("embedding", chunk_dict.get("chunk_id", ""), str(e))
                if args.fail_fast:
                    raise

        logger.info(f"  Embedded batch {i//10 + 1}, total: {len(embeddings)}/{len(chunks_to_embed)}")
        if not use_mock_embedding:
            time.sleep(0.5)  # Rate limit for real Bedrock calls

    # Write to LanceDB
    logger.info(f"Writing {len(embeddings)} vectors to LanceDB at {lancedb_path}")
    store = create_vector_store(
        backend="lancedb",
        db_path=lancedb_path,
        collection=collection_name,
    )

    # Build VectorStoreRecords
    records = []
    for emb in embeddings:
        records.append(VectorStoreRecord(
            chunk_id=emb["chunk_id"],
            document_id=emb.get("document_id", ""),
            text=emb.get("content", ""),
            embedding=emb.get("embedding", []),
            source_uri=emb.get("source_uri", ""),
            source_type=emb.get("source_type", ""),
            page=emb.get("page"),
            section_title=emb.get("section_title", ""),
            visual_block_ids=emb.get("visual_block_ids", []),
            acl=emb.get("acl", []),
            content_hash=emb.get("content_hash", ""),
            embedding_model=emb.get("embedding_model", "amazon.titan-embed-text-v2:0"),
            metadata={"run_id": state.run_id, "dataset": "murata"},
        ))

    if records:
        store.upsert_chunks(records)

    # Report
    report = {
        "run_id": state.run_id,
        "total_chunks": len(chunk_dicts),
        "embedded_count": len(embeddings),
        "lancedb_path": lancedb_path,
        "collection": collection_name,
        "records_written": len(records),
        "embedding_model": "amazon.titan-embed-text-v2:0",
        "dimension": 1024,
        "timestamp": datetime.now(timezone.utc).isoformat(),
    }
    state.save_json("lancedb_load_report.json", report)

    logger.info(f"Embedding complete: {len(embeddings)} vectors → LanceDB [{collection_name}]")
    return embeddings


# ===========================================================================
# Stage 5: GRAPH EXTRACTION
# ===========================================================================


def _mock_extraction(doc_chunk, run_id: str):
    """Generate mock extraction results for dry-run testing."""
    import hashlib
    from hermes_bedrock_agent.graph.extractor import ExtractionResult
    from hermes_bedrock_agent.schemas.graph import GraphEntity, GraphRelation

    # Generate a deterministic mock entity from chunk content
    name_words = doc_chunk.content[:60].split()[:3]
    entity_name = "_".join(name_words) if name_words else "MockEntity"
    eid = hashlib.sha256(f"{doc_chunk.chunk_id}:mock_ent".encode()).hexdigest()[:16]
    rid = hashlib.sha256(f"{doc_chunk.chunk_id}:mock_rel".encode()).hexdigest()[:16]

    entity = GraphEntity(
        entity_id=eid,
        canonical_name=entity_name,
        name=entity_name,
        entity_type="MockType",
        source_chunk_ids=[doc_chunk.chunk_id],
        confidence=0.5,
    )
    relation = GraphRelation(
        relation_id=rid,
        source_entity_id=eid,
        target_entity_id=eid,
        relation_type="MOCK_REL",
        source_chunk_ids=[doc_chunk.chunk_id],
        confidence=0.5,
    )
    return ExtractionResult(entities=[entity], relations=[relation], evidence=[], chunk_id=doc_chunk.chunk_id, errors=[])


def stage_graph(state: PipelineState, args) -> dict:
    """Extract entities/relations from chunks via Bedrock Claude."""
    logger.info("=" * 60)
    logger.info("STAGE 5: GRAPH — Entity/Relation extraction + normalization + QA")
    logger.info("=" * 60)

    from hermes_bedrock_agent.graph.extractor import GraphExtractor, ExtractorConfig, ExtractionResult
    from hermes_bedrock_agent.graph.normalizer import EntityNormalizer, NormalizerConfig
    from hermes_bedrock_agent.graph.quality_review import GraphQualityReviewer, QualityConfig
    from hermes_bedrock_agent.schemas.chunk import DocumentChunk, ChunkType
    from hermes_bedrock_agent.schemas.graph import GraphEntity, GraphRelation, EvidenceRecord

    # Load chunks
    chunk_dicts = state.load_jsonl("chunks.jsonl")
    if not chunk_dicts:
        raise RuntimeError("No chunks.jsonl found. Run chunk stage first.")

    # Apply exclude patterns for graph extraction
    docs_by_id = {}
    doc_dicts = state.load_jsonl("normalized_documents.jsonl")
    if doc_dicts:
        docs_by_id = {d["document_id"]: d for d in doc_dicts}

    exclude_patterns = getattr(args, "exclude_pattern", []) or []
    skip_graph_patterns = getattr(args, "skip_graph_extraction_pattern", []) or []
    chunk_dicts, filter_summary = _filter_chunks_for_stage(
        chunk_dicts, docs_by_id, "graph_extraction", exclude_patterns, skip_graph_patterns
    )
    state.save_jsonl("graph_extraction_filter_summary.json", [filter_summary])

    # Skip-existing: skip chunks already extracted
    existing_chunk_ids = set()
    if args.skip_existing or args.resume:
        existing_raw = state.load_jsonl("raw_entities.jsonl")
        existing_chunk_ids = {e.get("source_chunk_ids", [None])[0] for e in existing_raw if e.get("source_chunk_ids")}
        if existing_chunk_ids:
            logger.info(f"Skip-existing: {len(existing_chunk_ids)} chunks already extracted")

    # Filter: only extract from text/code chunks (not tiny ones)
    chunks_to_extract = [
        c for c in chunk_dicts
        if len(c.get("content", "")) > 50
        and c.get("chunk_id", "") not in existing_chunk_ids
    ]
    logger.info(f"Extracting from {len(chunks_to_extract)}/{len(chunk_dicts)} chunks")

    use_mock_graph = getattr(args, "mock_graph_extraction", False)
    neptune_dataset = getattr(args, "neptune_dataset", "murata")
    s3_uri = getattr(args, "s3_uri", DEFAULT_S3_URI)

    # Initialize extractor (only if not mocking)
    extractor = None
    if not use_mock_graph:
        extractor = GraphExtractor(config=ExtractorConfig(
            model_id="apac.anthropic.claude-sonnet-4-20250514-v1:0",
            max_tokens=4096,
            temperature=0.0,
            max_entities_per_chunk=8,
            max_relations_per_chunk=12,
            min_confidence=0.75,
        ))

    raw_entities = []
    raw_relations = []
    raw_evidence = []

    for i, chunk_dict in enumerate(chunks_to_extract):
        try:
            chunk_type_str = chunk_dict.get("chunk_type", "text")
            try:
                chunk_type = ChunkType(chunk_type_str)
            except (ValueError, KeyError):
                chunk_type = ChunkType.TEXT

            doc_chunk = DocumentChunk(
                chunk_id=chunk_dict["chunk_id"],
                document_id=chunk_dict.get("document_id", ""),
                content=chunk_dict.get("content", ""),
                chunk_type=chunk_type,
                chunk_index=chunk_dict.get("chunk_index", 0),
                source_uri=chunk_dict.get("source_uri", ""),
                source_type=chunk_dict.get("source_type", "unknown"),
                section_title=chunk_dict.get("section_title", ""),
                page=chunk_dict.get("page"),
                content_hash=chunk_dict.get("content_hash", ""),
            )

            result = extractor.extract_chunk(doc_chunk) if not use_mock_graph else _mock_extraction(doc_chunk, state.run_id)

            for ent in result.entities:
                ent_dict = ent.model_dump(mode="json")
                ent_dict["run_id"] = state.run_id
                ent_dict["dataset"] = neptune_dataset
                ent_dict["source_prefix"] = s3_uri
                raw_entities.append(ent_dict)
                state.append_jsonl("raw_entities.jsonl", ent_dict)

            for rel in result.relations:
                rel_dict = rel.model_dump(mode="json")
                rel_dict["run_id"] = state.run_id
                rel_dict["dataset"] = neptune_dataset
                rel_dict["source_prefix"] = s3_uri
                raw_relations.append(rel_dict)
                state.append_jsonl("raw_relations.jsonl", rel_dict)

            for ev in result.evidence:
                ev_dict = ev.model_dump(mode="json")
                ev_dict["run_id"] = state.run_id
                raw_evidence.append(ev_dict)
                state.append_jsonl("raw_evidence.jsonl", ev_dict)

            if (i + 1) % 5 == 0:
                logger.info(f"  Extracted {i+1}/{len(chunks_to_extract)}: "
                            f"{len(raw_entities)} entities, {len(raw_relations)} relations")
                time.sleep(1)  # Rate limit

        except Exception as e:
            logger.error(f"Extraction failed for chunk {chunk_dict.get('chunk_id', '?')}: {e}")
            state.record_failure("graph", chunk_dict.get("chunk_id", ""), str(e))
            if args.fail_fast:
                raise

    logger.info(f"Raw extraction: {len(raw_entities)} entities, {len(raw_relations)} relations")

    # Normalization
    logger.info("Running entity normalization...")
    normalizer = EntityNormalizer(config=NormalizerConfig())
    entity_models = []
    for e in raw_entities:
        try:
            from hermes_bedrock_agent.schemas.graph import EntityType
            entity_models.append(GraphEntity(
                entity_id=e["entity_id"],
                name=e.get("name", ""),
                canonical_name=e.get("canonical_name", ""),
                entity_type=EntityType(e.get("entity_type", "concept")),
                description=e.get("description", ""),
                aliases=e.get("aliases", []),
                source_chunk_ids=e.get("source_chunk_ids", []),
                confidence=e.get("confidence", 0.5),
                model_name=e.get("model_name", ""),
            ))
        except Exception:
            pass

    # Normalize each entity then deduplicate
    normalized_list = [normalizer.normalize_entity(e) for e in entity_models]
    normalized_entities = normalizer.deduplicate_entities(normalized_list)
    logger.info(f"Normalized: {len(entity_models)} → {len(normalized_entities)} entities")

    # Quality review
    logger.info("Running graph quality review...")
    from pathlib import Path as P
    schema_path = P(__file__).resolve().parent.parent / "configs" / "graph_schema.yaml"
    reviewer = GraphQualityReviewer(config=QualityConfig(min_confidence=0.75))

    relation_models = []
    for r in raw_relations:
        try:
            from hermes_bedrock_agent.schemas.graph import RelationType
            relation_models.append(GraphRelation(
                relation_id=r["relation_id"],
                source_entity_id=r.get("source_entity_id", ""),
                target_entity_id=r.get("target_entity_id", ""),
                relation_type=RelationType(r.get("relation_type", "related_to")),
                description=r.get("description", ""),
                source_chunk_id=r.get("source_chunk_id", ""),
                source_chunk_ids=r.get("source_chunk_ids", []),
                confidence=r.get("confidence", 0.5),
                weight=r.get("weight", 1.0),
                evidence_id=r.get("evidence_id", ""),
            ))
        except Exception:
            pass

    evidence_models = []
    for ev in raw_evidence:
        try:
            evidence_models.append(EvidenceRecord(
                evidence_id=ev.get("evidence_id", ""),
                entity_id=ev.get("entity_id"),
                relation_id=ev.get("relation_id"),
                source_chunk_id=ev.get("source_chunk_id", ev.get("chunk_id", "")),
                document_id=ev.get("document_id", ""),
                source_uri=ev.get("source_uri", ""),
                evidence_text=ev.get("evidence_text", ev.get("text", ev.get("content", ""))),
                context_text=ev.get("context_text", ""),
                page=ev.get("page"),
                section_title=ev.get("section_title", ""),
                confidence=ev.get("confidence", 0.5),
                model_name=ev.get("model_name", ""),
            ))
        except Exception:
            pass

    review = reviewer.review(normalized_entities, relation_models, evidence_models)

    # Save normalized/reviewed results
    accepted_entities = [e.model_dump(mode="json") for e in review.accepted_entities]
    accepted_relations = [r.model_dump(mode="json") for r in review.accepted_relations]
    pending_relations = [r.model_dump(mode="json") for r in review.pending_relations]
    rejected_relations = [r.model_dump(mode="json") for r in review.rejected_relations]

    state.save_jsonl("entities.jsonl", accepted_entities)
    state.save_jsonl("relations.jsonl", accepted_relations)
    state.save_jsonl("evidence.jsonl", raw_evidence)
    state.save_jsonl("pending_relations.jsonl", pending_relations)
    state.save_jsonl("rejected_relations.jsonl", rejected_relations)

    # Quality report
    quality_report = {
        "run_id": state.run_id,
        "raw_entities": len(raw_entities),
        "raw_relations": len(raw_relations),
        "normalized_entities": len(normalized_entities),
        "accepted_entities": len(accepted_entities),
        "accepted_relations": len(accepted_relations),
        "pending_relations": len(pending_relations),
        "rejected_relations": len(rejected_relations),
        "evidence_records": len(raw_evidence),
        "timestamp": datetime.now(timezone.utc).isoformat(),
    }
    state.save_json("graph_quality_report.json", quality_report)

    logger.info(f"Graph stage complete: {len(accepted_entities)} entities, "
                f"{len(accepted_relations)} relations accepted")
    return quality_report


# ===========================================================================
# Stage 6: NEPTUNE LOAD
# ===========================================================================


def stage_load(state: PipelineState, args) -> dict:
    """Load graph to Neptune Analytics using parameterized queries."""
    logger.info("=" * 60)
    logger.info("STAGE 6: LOAD — Neptune Graph write (parameterized)")
    logger.info("=" * 60)

    from hermes_bedrock_agent.clients.neptune_client import NeptuneClient
    from hermes_bedrock_agent.graph.neptune_loader import (
        build_import_cypher,
        load_to_neptune,
        write_neptune_import_cypher,
    )
    from hermes_bedrock_agent.schemas.graph import EntityType, GraphEntity, GraphRelation, RelationType

    # Load entities/relations
    entity_dicts = state.load_jsonl("entities.jsonl")
    relation_dicts = state.load_jsonl("relations.jsonl")

    if not entity_dicts:
        raise RuntimeError("No entities.jsonl found. Run graph stage first.")

    # Rebuild models
    entities = []
    for e in entity_dicts:
        try:
            ent = GraphEntity(
                entity_id=e["entity_id"],
                name=e.get("name", ""),
                canonical_name=e.get("canonical_name", ""),
                entity_type=EntityType(e.get("entity_type", "concept")),
                description=e.get("description", ""),
                aliases=e.get("aliases", []),
                source_chunk_ids=e.get("source_chunk_ids", []),
                confidence=e.get("confidence", 0.5),
                model_name=e.get("model_name", ""),
                acl=e.get("acl", []),
            )
            entities.append(ent)
        except Exception as ex:
            logger.warning(f"Skipping entity {e.get('entity_id', '?')}: {ex}")

    relations = []
    for r in relation_dicts:
        try:
            rel = GraphRelation(
                relation_id=r["relation_id"],
                source_entity_id=r.get("source_entity_id", ""),
                target_entity_id=r.get("target_entity_id", ""),
                relation_type=RelationType(r.get("relation_type", "related_to")),
                description=r.get("description", ""),
                source_chunk_id=r.get("source_chunk_id", ""),
                source_chunk_ids=r.get("source_chunk_ids", []),
                confidence=r.get("confidence", 0.5),
                weight=r.get("weight", 1.0),
                evidence_id=r.get("evidence_id", ""),
                acl=r.get("acl", []),
            )
            relations.append(rel)
        except Exception as ex:
            logger.warning(f"Skipping relation {r.get('relation_id', '?')}: {ex}")

    logger.info(f"Loading {len(entities)} entities + {len(relations)} relations")

    # ALWAYS export inline Cypher (for review/artifact)
    cypher_path = state.artifact_path("neptune_import.cypher")
    write_neptune_import_cypher(entities, relations, cypher_path)

    # Live Neptune write
    if args.live_neptune and args.confirm_live_write:
        logger.info(f"LIVE WRITE to Neptune: {args.neptune_endpoint}")

        # Extract graph_id from endpoint
        # Endpoint format: g-XXXXX.region.neptune-graph.amazonaws.com
        graph_id = args.neptune_endpoint.split(".")[0]

        neptune = NeptuneClient(
            graph_id=graph_id,
            region="ap-northeast-1",
        )

        # Inject run_id and dataset into all entity/relation properties
        neptune_dataset = getattr(args, "neptune_dataset", "murata")
        s3_uri = getattr(args, "s3_uri", DEFAULT_S3_URI)
        for ent in entities:
            if ent.metadata is None:
                ent.metadata = {}
            ent.metadata["run_id"] = state.run_id
            ent.metadata["dataset"] = neptune_dataset
            ent.metadata["source_prefix"] = s3_uri
        for rel in relations:
            if rel.metadata is None:
                rel.metadata = {}
            rel.metadata["run_id"] = state.run_id
            rel.metadata["dataset"] = neptune_dataset
            rel.metadata["source_prefix"] = s3_uri

        result = load_to_neptune(
            neptune,
            entities,
            relations,
            dry_run=False,
            batch_size=20,
        )

        # Verification query
        verify_result = {}
        try:
            count_result = neptune.execute_query(
                "MATCH (n) WHERE n.entity_id IS NOT NULL RETURN count(n) AS node_count"
            )
            verify_result["node_count"] = count_result
        except Exception as e:
            verify_result["error"] = str(e)

        report = {
            "run_id": state.run_id,
            "neptune_endpoint": args.neptune_endpoint,
            "graph_id": graph_id,
            "nodes_loaded": result["nodes_loaded"],
            "edges_loaded": result["edges_loaded"],
            "errors": result["errors"],
            "dry_run": False,
            "verification": verify_result,
            "timestamp": datetime.now(timezone.utc).isoformat(),
        }
    else:
        logger.info("Neptune write SKIPPED (use --live-neptune --confirm-live-write)")
        report = {
            "run_id": state.run_id,
            "neptune_endpoint": args.neptune_endpoint,
            "nodes_to_load": len(entities),
            "edges_to_load": len(relations),
            "dry_run": True,
            "cypher_export": str(cypher_path),
            "timestamp": datetime.now(timezone.utc).isoformat(),
        }

    state.save_json("neptune_load_report.json", report)

    # Generate verification doc
    verification_md = _generate_load_verification(state, report, entities, relations)
    state.save_text("load_verification_report.md", verification_md)

    logger.info(f"Load stage complete: {report}")
    return report


def _generate_load_verification(state, report, entities, relations) -> str:
    """Generate load verification markdown report."""
    lines = [
        f"# Neptune Load Verification Report",
        f"",
        f"**Run ID:** {state.run_id}",
        f"**Endpoint:** {report.get('neptune_endpoint', 'N/A')}",
        f"**Timestamp:** {report.get('timestamp', '')}",
        f"",
        f"## Summary",
        f"",
        f"| Metric | Value |",
        f"|--------|-------|",
        f"| Entities to load | {len(entities)} |",
        f"| Relations to load | {len(relations)} |",
        f"| Nodes loaded | {report.get('nodes_loaded', 'N/A')} |",
        f"| Edges loaded | {report.get('edges_loaded', 'N/A')} |",
        f"| Errors | {report.get('errors', 'N/A')} |",
        f"| Dry run | {report.get('dry_run', True)} |",
        f"",
        f"## Entity Type Distribution",
        f"",
    ]

    type_counts: dict[str, int] = {}
    for e in entities:
        t = e.entity_type.value
        type_counts[t] = type_counts.get(t, 0) + 1
    for t, c in sorted(type_counts.items(), key=lambda x: -x[1]):
        lines.append(f"- {t}: {c}")

    lines.extend([
        f"",
        f"## Verification Queries",
        f"",
        f"```cypher",
        f"-- Count all nodes",
        f"MATCH (n) RETURN count(n) AS total_nodes",
        f"",
        f"-- Count by label",
        f"MATCH (n) RETURN labels(n)[0] AS label, count(n) AS cnt ORDER BY cnt DESC",
        f"",
        f"-- Count edges",
        f"MATCH ()-[r]->() RETURN count(r) AS total_edges",
        f"",
        f"-- Sample nodes",
        f"MATCH (n) RETURN n LIMIT 5",
        f"```",
    ])

    return "\n".join(lines)


# ===========================================================================
# ===========================================================================
# Stage 7: ENRICHMENT (Optional — skipped by default)
# ===========================================================================


def stage_enrichment(state: PipelineState, args) -> None:
    """Optional i18n enrichment stage — SKIPPED unless explicitly enabled.

    Adds multilingual display names, business aliases, and zh/en/ja labels.
    Default mode is 'none' which skips entirely. No LLM calls, no Neptune writes.
    """
    logger.info("=" * 60)
    logger.info("STAGE 7: ENRICHMENT (Optional)")
    logger.info(f"  Mode: {args.enrichment_mode}")
    logger.info("=" * 60)

    mode = args.enrichment_mode

    if mode == "none":
        logger.info("  ⏭ Enrichment SKIPPED (mode=none, default behavior)")
        logger.info("  To enable: --enrichment-mode rule|mock|llm")
        return

    # Safety check: LLM mode requires live-source or full-live pipeline mode
    if mode == "llm" and args.mode == "mock-dry-run":
        logger.warning("  ⚠ enrichment-mode=llm requires --mode live-source or full-live")
        logger.warning("  Falling back to enrichment-mode=mock")
        mode = "mock"

    # Safety check: Neptune write requires confirm
    if args.enrichment_update_neptune and not args.confirm_live_write:
        logger.error("  ✗ --enrichment-update-neptune requires --confirm-live-write")
        logger.error("  Skipping Neptune write, generating preview only")
        args.enrichment_update_neptune = False

    from hermes_bedrock_agent.graph.i18n_enricher import (
        run_enrichment,
    )

    entities_path = state.artifact_path("entities.jsonl")
    relations_clean_path = state.artifact_path("relations_clean.jsonl")

    if not entities_path.exists():
        logger.error(f"  ✗ entities.jsonl not found at {entities_path}")
        return

    # Load entities
    entities = state.load_jsonl("entities.jsonl")
    relations = state.load_jsonl("relations_clean.jsonl") if relations_clean_path.exists() else []

    logger.info(f"  Entities: {len(entities)}")
    logger.info(f"  Relations: {len(relations)}")
    logger.info(f"  Max entities: {args.enrichment_max_entities}")
    logger.info(f"  Update Neptune: {args.enrichment_update_neptune}")

    result = run_enrichment(
        mode=mode,
        entities=entities,
        relations=relations,
        max_entities=args.enrichment_max_entities,
        output_dir=state.artifact_dir,
        output_suffix=args.enrichment_output_suffix,
        update_neptune=args.enrichment_update_neptune,
    )

    if result:
        logger.info(f"  ✓ Enrichment complete (mode={mode})")
        logger.info(f"    Entities enriched: {result.get('entities_enriched', 0)}")
        logger.info(f"    Relations enriched: {result.get('relations_enriched', 0)}")
        if result.get("output_files"):
            for f in result["output_files"]:
                logger.info(f"    Output: {f}")
    else:
        logger.info("  ⏭ Enrichment produced no output")


# ===========================================================================
# Stage 8: RETRIEVAL
# ===========================================================================


def stage_retrieval(state: PipelineState, args) -> None:
    """Run retrieval demo: LanceDB vector + Neptune graph → fused answers."""
    logger.info("=" * 60)
    logger.info("STAGE 8: RETRIEVAL — Hybrid search + answer generation")
    logger.info("=" * 60)

    from hermes_bedrock_agent.embedding.embedder import BedrockEmbedder, EmbedderConfig
    from hermes_bedrock_agent.retrieval.text_retriever import VectorStoreTextRetriever, TextRetrieverConfig
    from hermes_bedrock_agent.retrieval.fusion import fuse_evidence, FusionConfig, FusionStrategy
    from hermes_bedrock_agent.generation.answer_generator import AnswerGenerator, AnswerGeneratorConfig
    from hermes_bedrock_agent.vector_store import create_vector_store

    # Demo queries for Murata project
    demo_queries = [
        "仕訳基礎テーブルの構造を説明してください",
        "付款申請の承認フローはどのように動作しますか",
        "PaymentReqAction.javaの主要機能は何ですか",
        "対帳単（Receiving List）の処理ロジックを説明してください",
        "システム管理画面にはどのような機能がありますか",
    ]

    # Initialize components
    lancedb_path = str(getattr(args, "local_vector_store_path", None) or DEFAULT_LANCEDB_PATH)
    collection_name = getattr(args, "local_vector_collection", None) or f"murata_e2e_{state.run_id}"

    store = create_vector_store(
        backend="lancedb",
        db_path=lancedb_path,
        collection=collection_name,
    )

    retriever = VectorStoreTextRetriever(
        store,
        config=TextRetrieverConfig(top_k=10),
    )

    embedder = BedrockEmbedder(config=EmbedderConfig(
        model_id="amazon.titan-embed-text-v2:0",
        dimension=1024,
    ))

    # Neptune graph retriever (if live)
    graph_retriever = None
    if args.live_neptune:
        try:
            from hermes_bedrock_agent.clients.neptune_client import NeptuneClient
            from hermes_bedrock_agent.retrieval.graph_retriever import NeptuneGraphRetriever

            graph_id = args.neptune_endpoint.split(".")[0]
            neptune = NeptuneClient(graph_id=graph_id, region="ap-northeast-1")
            graph_retriever = NeptuneGraphRetriever(neptune)
        except Exception as e:
            logger.warning(f"Neptune graph retriever init failed: {e}")

    # Answer generator
    use_mock_answer = getattr(args, "mock_answer", False)
    answer_gen = None
    if not use_mock_answer:
        import boto3
        bedrock_runtime = boto3.client("bedrock-runtime", region_name="ap-northeast-1")
        answer_gen = AnswerGenerator(
            bedrock_client=bedrock_runtime,
            config=AnswerGeneratorConfig(
                model_id="apac.anthropic.claude-sonnet-4-20250514-v1:0",
                max_tokens=2048,
            ),
        )

    fused_contexts = []
    answer_results = []
    retrieval_examples = []

    for query in demo_queries:
        try:
            logger.info(f"Query: {query}")

            # 1. Embed query
            query_embedding = embedder.embed_text(query)

            # 2. Text retrieval via LanceDB
            text_evidence = retriever.hybrid_search(query, query_embedding=query_embedding)
            logger.info(f"  Text evidence: {len(text_evidence)} results")

            # 3. Graph retrieval (if available)
            graph_evidence = []
            if graph_retriever:
                try:
                    # Extract key terms for graph search
                    terms = [t for t in query.split() if len(t) > 2][:5]
                    graph_evidence = graph_retriever.retrieve_graph_context(terms)
                    logger.info(f"  Graph evidence: {len(graph_evidence)} results")
                except Exception as e:
                    logger.warning(f"  Graph retrieval failed: {e}")

            # 4. Fuse evidence
            fused = fuse_evidence(
                text_evidence, graph_evidence,
                query=query,
                config=FusionConfig(strategy=FusionStrategy.RRF),
            )
            fused_dict = fused.model_dump(mode="json")
            fused_dict["query"] = query
            fused_contexts.append(fused_dict)

            # 5. Generate answer
            try:
                if use_mock_answer:
                    answer_dict = {
                        "query": query,
                        "answer": f"[MOCK] Answer for: {query[:100]}",
                        "citations": [],
                        "model_name": "mock",
                        "confidence": 0.5,
                    }
                    answer_results.append(answer_dict)
                    logger.info(f"  Answer (mock): {answer_dict['answer'][:80]}")
                else:
                    answer = answer_gen.generate_answer(query, fused)
                    answer_dict = answer.model_dump(mode="json")
                    answer_dict["query"] = query
                    answer_results.append(answer_dict)
                    logger.info(f"  Answer: {answer.answer[:100]}...")
            except Exception as e:
                logger.warning(f"  Answer generation failed: {e}")
                answer_results.append({"query": query, "error": str(e)})

            # Retrieval example
            retrieval_examples.append({
                "query": query,
                "text_evidence_count": len(text_evidence),
                "graph_evidence_count": len(graph_evidence),
                "fused_evidence_count": fused.total_evidence_count,
                "top_sources": [e.source_uri for e in text_evidence[:3]],
            })

            time.sleep(1)  # Rate limit

        except Exception as e:
            logger.error(f"Retrieval failed for query '{query}': {e}")
            state.record_failure("retrieval", query, str(e))
            if args.fail_fast:
                raise

    # Save artifacts
    state.save_jsonl("fused_context_examples.jsonl", fused_contexts)
    state.save_jsonl("answer_examples.jsonl", answer_results)
    state.save_jsonl("retrieval_live_examples.jsonl", retrieval_examples)

    logger.info(f"Retrieval stage complete: {len(answer_results)} answers generated")


# ===========================================================================
# Stage 8: VISUALIZATION
# ===========================================================================


def stage_visualization(state: PipelineState, args) -> None:
    """Generate Mermaid diagrams from Neptune subgraphs."""
    logger.info("=" * 60)
    logger.info("STAGE 9: VISUALIZATION — Mermaid diagram generation")
    logger.info("=" * 60)

    from hermes_bedrock_agent.visualization.mermaid_generator import MermaidGenerator, MermaidConfig
    from hermes_bedrock_agent.schemas.visualization import SubgraphResult, VisualizationNode, VisualizationEdge

    # Load entities/relations for local visualization
    entity_dicts = state.load_jsonl("entities.jsonl")
    relation_dicts = state.load_jsonl("relations.jsonl")

    if not entity_dicts:
        logger.warning("No entities found, skipping visualization")
        state.save_text("mermaid_examples.md", "# Mermaid Examples\n\nNo entities found.\n")
        return

    # Build local subgraph for visualization
    nodes = []
    for e in entity_dicts[:50]:  # Limit for readability
        nodes.append(VisualizationNode(
            node_id=e["entity_id"],
            label=e.get("name", e.get("canonical_name", "?")),
            node_type=e.get("entity_type", "unknown"),
            properties={"description": e.get("description", "")},
        ))

    edges = []
    node_ids = {n.node_id for n in nodes}
    for r in relation_dicts:
        if r.get("source_entity_id") in node_ids and r.get("target_entity_id") in node_ids:
            edges.append(VisualizationEdge(
                edge_id=r.get("relation_id", ""),
                source_id=r["source_entity_id"],
                target_id=r["target_entity_id"],
                edge_type=r.get("relation_type", "related_to"),
                label=r.get("description", "")[:30],
            ))

    subgraph = SubgraphResult(
        nodes=nodes,
        edges=edges,
        center_entity="overview",
        depth=2,
        query_time_ms=0,
    )

    # Generate Mermaid
    generator = MermaidGenerator(config=MermaidConfig(
        direction="LR",
        max_nodes=30,
        show_edge_labels=True,
    ))

    mermaid_code = generator.generate(subgraph)

    # Also try Neptune live query if available
    neptune_mermaid = ""
    if args.live_neptune:
        try:
            from hermes_bedrock_agent.clients.neptune_client import NeptuneClient
            from hermes_bedrock_agent.visualization.subgraph_query import SubgraphQueryService

            graph_id = args.neptune_endpoint.split(".")[0]
            neptune = NeptuneClient(graph_id=graph_id, region="ap-northeast-1")
            query_svc = SubgraphQueryService(neptune)

            # Pick a central entity
            if entity_dicts:
                center = entity_dicts[0].get("entity_id", "")
                live_subgraph = query_svc.query_subgraph(center, depth=2, max_nodes=30)
                neptune_mermaid = generator.generate(live_subgraph)
        except Exception as e:
            logger.warning(f"Neptune visualization failed: {e}")
            neptune_mermaid = f"<!-- Neptune query failed: {e} -->"

    # Save Mermaid markdown
    md_lines = [
        "# Mermaid Visualization Examples",
        "",
        f"**Run ID:** {state.run_id}",
        f"**Entities:** {len(entity_dicts)}",
        f"**Relations:** {len(relation_dicts)}",
        "",
        "## Local Subgraph (from artifacts)",
        "",
        "```mermaid",
        mermaid_code,
        "```",
        "",
    ]

    if neptune_mermaid:
        md_lines.extend([
            "## Neptune Live Subgraph",
            "",
            "```mermaid",
            neptune_mermaid,
            "```",
            "",
        ])

    state.save_text("mermaid_examples.md", "\n".join(md_lines))
    logger.info("Visualization stage complete")


# ===========================================================================
# Final Reports
# ===========================================================================


def generate_final_reports(state: PipelineState, args) -> None:
    """Generate quality report and cleanup commands."""
    logger.info("Generating final reports...")

    elapsed = time.time() - state.start_time

    # Quality report
    quality_lines = [
        "# Murata E2E Quality Report",
        "",
        f"**Run ID:** {state.run_id}",
        f"**Dataset:** murata",
        f"**Source:** {DEFAULT_S3_URI}",
        f"**Duration:** {elapsed:.1f}s",
        f"**Timestamp:** {datetime.now(timezone.utc).isoformat()}",
        "",
        "## Pipeline Summary",
        "",
    ]

    # Load stats from each stage
    inventory = _load_json_safe(state.artifact_path("file_inventory_report.json"))
    chunk_stats = _load_json_safe(state.artifact_path("chunk_stats.json"))
    lancedb_report = _load_json_safe(state.artifact_path("lancedb_load_report.json"))
    graph_report = _load_json_safe(state.artifact_path("graph_quality_report.json"))
    neptune_report = _load_json_safe(state.artifact_path("neptune_load_report.json"))

    quality_lines.extend([
        "| Stage | Metric | Value |",
        "|-------|--------|-------|",
        f"| Scan | Documents | {inventory.get('processable_documents', 'N/A')} |",
        f"| Scan | Total bytes | {inventory.get('total_bytes', 'N/A')} |",
        f"| Chunk | Total chunks | {chunk_stats.get('total_chunks', 'N/A')} |",
        f"| Embedding | Vectors | {lancedb_report.get('embedded_count', 'N/A')} |",
        f"| Graph | Entities | {graph_report.get('accepted_entities', 'N/A')} |",
        f"| Graph | Relations | {graph_report.get('accepted_relations', 'N/A')} |",
        f"| Neptune | Nodes loaded | {neptune_report.get('nodes_loaded', 'N/A')} |",
        f"| Neptune | Edges loaded | {neptune_report.get('edges_loaded', 'N/A')} |",
        "",
    ])

    # Failures
    if state.failures:
        quality_lines.extend([
            "## Failures",
            "",
            f"Total failures: {len(state.failures)}",
            "",
        ])
        for f in state.failures[:20]:
            quality_lines.append(f"- [{f['stage']}] {f['item_id']}: {f['error'][:80]}")
        quality_lines.append("")

    state.save_text("murata_e2e_quality_report.md", "\n".join(quality_lines))

    # Cleanup commands
    collection_name = getattr(args, "local_vector_collection", None) or f"murata_e2e_{state.run_id}"
    lancedb_path = str(getattr(args, "local_vector_store_path", None) or DEFAULT_LANCEDB_PATH)
    neptune_dataset = getattr(args, "neptune_dataset", "murata")
    neptune_endpoint = getattr(args, "neptune_endpoint", DEFAULT_NEPTUNE_ENDPOINT)

    cleanup_lines = [
        "# Cleanup Commands",
        "",
        f"**Run ID:** {state.run_id}",
        f"**Dataset:** {neptune_dataset}",
        f"**Collection:** {collection_name}",
        f"**Neptune Endpoint:** {neptune_endpoint}",
        "",
        "## 1. Delete LanceDB Collection",
        "",
        "```python",
        "import lancedb",
        f'db = lancedb.connect("{lancedb_path}")',
        f'db.drop_table("{collection_name}")',
        f'print("Dropped LanceDB collection: {collection_name}")',
        "```",
        "",
        "## 2. Delete Neptune Test Data (by run_id — parameterized)",
        "",
        "**WARNING:** Review before executing. Only use on test graphs.",
        "",
        "```python",
        "# Use parameterized openCypher — no string interpolation",
        "import boto3",
        "from botocore.auth import SigV4Auth",
        "",
        f'NEPTUNE_ENDPOINT = "{neptune_endpoint}"',
        f'RUN_ID = "{state.run_id}"',
        f'DATASET = "{neptune_dataset}"',
        "",
        "# Step 1: Delete relations by run_id",
        "query_del_rels = '''",
        "MATCH ()-[r]->() WHERE r.run_id = $run_id",
        "DELETE r",
        "RETURN count(r) AS deleted_rels",
        "'''",
        "",
        "# Step 2: Delete entities by run_id",
        "query_del_nodes = '''",
        "MATCH (n) WHERE n.run_id = $run_id",
        "DETACH DELETE n",
        "RETURN count(n) AS deleted_nodes",
        "'''",
        "",
        "# Execute via NeptuneClient (parameterized):",
        "from hermes_bedrock_agent.clients.neptune_client import NeptuneClient",
        f'client = NeptuneClient(endpoint="{neptune_endpoint}")',
        'result1 = client.execute_parameterized_query(query_del_rels, {"run_id": RUN_ID})',
        'result2 = client.execute_parameterized_query(query_del_nodes, {"run_id": RUN_ID})',
        'print(f"Deleted rels: {result1}, nodes: {result2}")',
        "```",
        "",
        "Alternative (filter by dataset):",
        "```cypher",
        f"MATCH (n) WHERE n.dataset = $dataset AND n.run_id = $run_id DETACH DELETE n",
        f"// params: {{dataset: '{neptune_dataset}', run_id: '{state.run_id}'}}",
        "```",
        "",
        "## 3. Delete Local Artifacts",
        "",
        "```bash",
        f"rm -rf {state.artifact_dir}",
        f'echo "Deleted artifacts: {state.artifact_dir}"',
        "```",
        "",
        "## 4. Delete All Run Data (one-shot)",
        "",
        "```bash",
        f"# Delete artifacts",
        f"rm -rf {state.artifact_dir}",
        f"",
        f"# Delete LanceDB collection",
        f'python -c "import lancedb; db = lancedb.connect(\'{lancedb_path}\'); db.drop_table(\'{collection_name}\')"',
        f"",
        f"# Verify cleanup",
        f'python -c "import lancedb; db = lancedb.connect(\'{lancedb_path}\'); print(db.table_names())"',
        "```",
    ]

    state.save_text("cleanup_commands.md", "\n".join(cleanup_lines))
    logger.info("Final reports generated")


def _load_json_safe(path: Path) -> dict:
    """Load JSON file, return empty dict if not found."""
    if path.exists():
        with open(path) as f:
            return json.load(f)
    return {}


# ===========================================================================
# CLI
# ===========================================================================


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="E2E Murata Pipeline — Full live GraphRAG pipeline",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Modes:
  mock-dry-run      (default) — No AWS access. Parse from local cache or fail gracefully.
  live-source       --mode live-source — Access S3 + Bedrock but do NOT write Neptune.
  full-live         --mode full-live — Access S3 + Bedrock + LanceDB write + Neptune write.
                    Requires: --confirm-live-write --neptune-endpoint --run-id

Examples:
  # Mock dry-run (no AWS access):
  python scripts/run_e2e_murata_pipeline.py --mode mock-dry-run --stage scan

  # Live source (S3 + Bedrock reads, LanceDB writes, no Neptune):
  python scripts/run_e2e_murata_pipeline.py --mode live-source --run-id test_001

  # Full live (all services):
  python scripts/run_e2e_murata_pipeline.py --mode full-live \\
    --run-id murata_prod_001 --confirm-live-write \\
    --neptune-endpoint g-nbuyck5yl8.ap-northeast-1.neptune-graph.amazonaws.com
""",
    )

    # Stage selection
    parser.add_argument("--stage", default="all",
                        choices=STAGE_ORDER + ["all"],
                        help="Stage to run (default: all)")

    # Run identity
    parser.add_argument("--run-id", default=DEFAULT_RUN_ID,
                        help=f"Run identifier (default: {DEFAULT_RUN_ID})")
    parser.add_argument("--artifact-base", type=Path, default=DEFAULT_ARTIFACT_BASE,
                        help="Base directory for artifacts")

    # Execution mode
    parser.add_argument("--mode", default="mock-dry-run",
                        choices=["mock-dry-run", "live-source", "full-live"],
                        help="Execution mode (default: mock-dry-run)")

    # S3 source
    parser.add_argument("--s3-uri", default=DEFAULT_S3_URI,
                        help=f"S3 source prefix (default: {DEFAULT_S3_URI})")

    # VLM flags
    parser.add_argument("--enable-vlm", action="store_true", default=True,
                        help="Enable VLM parsing for images/PDF pages (default: True)")
    parser.add_argument("--no-vlm", action="store_true",
                        help="Disable VLM parsing")

    # Mock flags (for mock-dry-run mode)
    parser.add_argument("--mock-embedding", action="store_true",
                        help="Use mock embeddings (random vectors) instead of Bedrock")
    parser.add_argument("--mock-graph-extraction", action="store_true",
                        help="Use mock graph extraction instead of Bedrock Claude")
    parser.add_argument("--mock-answer", action="store_true",
                        help="Use mock answer generation instead of Bedrock Claude")

    # Neptune flags
    parser.add_argument("--live-neptune", action="store_true",
                        help="Enable live Neptune write (implied by --mode full-live)")
    parser.add_argument("--confirm-live-write", action="store_true",
                        help="Confirm live write to Neptune (safety flag, required for full-live)")
    parser.add_argument("--neptune-endpoint", default=DEFAULT_NEPTUNE_ENDPOINT,
                        help=f"Neptune endpoint (default: {DEFAULT_NEPTUNE_ENDPOINT})")
    parser.add_argument("--neptune-dataset", default="murata",
                        help="Neptune dataset tag written to all nodes/edges (default: murata)")

    # Resilience flags
    parser.add_argument("--resume", action="store_true",
                        help="Resume from checkpoint (reload existing artifacts)")
    parser.add_argument("--skip-existing", action="store_true",
                        help="Skip documents/chunks already processed")
    parser.add_argument("--fail-fast", action="store_true",
                        help="Halt on first error (default: continue)")
    parser.add_argument("--save-failures", action="store_true", default=True,
                        help="Save failures to parser_failed.jsonl")

    # Exclude patterns
    parser.add_argument("--exclude-pattern", action="append", default=[],
                        help="Glob pattern to exclude from both embedding and graph extraction. "
                             "Matches against filename and source_uri. Repeatable.")
    parser.add_argument("--skip-embedding-pattern", action="append", default=[],
                        help="Glob pattern to exclude from embedding only. "
                             "Matches against filename and source_uri. Repeatable.")
    parser.add_argument("--skip-graph-extraction-pattern", action="append", default=[],
                        help="Glob pattern to exclude from graph extraction only. "
                             "Matches against filename and source_uri. Repeatable.")

    # Vector store
    parser.add_argument("--vector-store-backend", default="lancedb",
                        help="Vector store backend (default: lancedb)")
    parser.add_argument("--local-vector-store-path", type=Path, default=DEFAULT_LANCEDB_PATH,
                        help=f"Local vector store path (default: {DEFAULT_LANCEDB_PATH})")
    parser.add_argument("--local-vector-collection", default=None,
                        help="LanceDB collection name (default: murata_e2e_{run_id})")

    # Legacy alias
    parser.add_argument("--lancedb-path", type=Path, default=None,
                        help=argparse.SUPPRESS)  # Hidden alias for --local-vector-store-path

    # --- Enrichment (Optional Stage — skipped by default) ---
    parser.add_argument("--enrichment-mode", default="none",
                        choices=["none", "rule", "mock", "llm"],
                        dest="enrichment_mode",
                        help="Enrichment mode (default: none — skip enrichment entirely)")
    parser.add_argument("--enrichment-max-entities", type=int, default=200,
                        dest="enrichment_max_entities",
                        help="Max entities to enrich (default: 200)")
    parser.add_argument("--enrichment-update-neptune", action="store_true", default=False,
                        dest="enrichment_update_neptune",
                        help="Write enrichment results to Neptune (requires --confirm-live-write)")
    parser.add_argument("--enrichment-output-suffix", default="",
                        dest="enrichment_output_suffix",
                        help="Suffix for enrichment output files")

    return parser.parse_args()


def main():
    args = parse_args()

    if args.no_vlm:
        args.enable_vlm = False

    # Handle legacy --lancedb-path alias
    if args.lancedb_path is not None:
        args.local_vector_store_path = args.lancedb_path

    # Derive LanceDB collection name from run_id
    if args.local_vector_collection is None:
        args.local_vector_collection = f"murata_e2e_{args.run_id}"

    # Mode-based flag propagation
    if args.mode == "mock-dry-run":
        args.mock_embedding = True
        args.mock_graph_extraction = True
        args.mock_answer = True
        args.live_neptune = False
    elif args.mode == "live-source":
        # Access S3 + Bedrock, write LanceDB, but NOT Neptune
        args.live_neptune = False
    elif args.mode == "full-live":
        args.live_neptune = True

    # Validation for full-live mode
    if args.mode == "full-live":
        errors = []
        if not args.confirm_live_write:
            errors.append("--confirm-live-write is required for full-live mode")
        if not args.neptune_endpoint:
            errors.append("--neptune-endpoint is required for full-live mode")
        if args.run_id == DEFAULT_RUN_ID:
            errors.append(
                "--run-id must be explicitly set for full-live mode "
                f"(default '{DEFAULT_RUN_ID}' is not allowed)"
            )
        if errors:
            logger.error("FULL-LIVE MODE VALIDATION FAILED:")
            for e in errors:
                logger.error(f"  ✗ {e}")
            sys.exit(1)

    # Set environment
    os.environ["VECTOR_STORE_BACKEND"] = args.vector_store_backend
    os.environ.setdefault("AWS_DEFAULT_REGION", "ap-northeast-1")

    logger.info("=" * 70)
    logger.info(f"  MURATA E2E PIPELINE — Run: {args.run_id}")
    logger.info(f"  Mode: {args.mode}")
    logger.info(f"  Stage: {args.stage}")
    logger.info(f"  S3 URI: {args.s3_uri}")
    logger.info(f"  VLM: {args.enable_vlm}")
    logger.info(f"  Mock embedding: {args.mock_embedding}")
    logger.info(f"  Mock graph extraction: {args.mock_graph_extraction}")
    logger.info(f"  Mock answer: {args.mock_answer}")
    logger.info(f"  Neptune: {args.live_neptune} (confirm: {args.confirm_live_write})")
    logger.info(f"  Neptune endpoint: {args.neptune_endpoint}")
    logger.info(f"  Neptune dataset: {args.neptune_dataset}")
    logger.info(f"  Enrichment: mode={args.enrichment_mode}, max={args.enrichment_max_entities}")
    logger.info(f"  Vector store: {args.vector_store_backend} @ {args.local_vector_store_path}")
    logger.info(f"  Collection: {args.local_vector_collection}")
    logger.info(f"  Resume: {args.resume}, Skip-existing: {args.skip_existing}")
    if args.exclude_pattern:
        logger.info(f"  Exclude patterns: {args.exclude_pattern}")
    if args.skip_embedding_pattern:
        logger.info(f"  Skip-embedding patterns: {args.skip_embedding_pattern}")
    if args.skip_graph_extraction_pattern:
        logger.info(f"  Skip-graph-extraction patterns: {args.skip_graph_extraction_pattern}")
    logger.info("=" * 70)

    state = PipelineState(args.run_id, args.artifact_base)

    stages_to_run = STAGE_ORDER if args.stage == "all" else [args.stage]

    for stage_name in stages_to_run:
        try:
            if stage_name == "scan":
                stage_scan(state, args)
            elif stage_name == "parse":
                stage_parse(state, args)
            elif stage_name == "chunk":
                stage_chunk(state, args)
            elif stage_name == "embedding":
                stage_embedding(state, args)
            elif stage_name == "graph":
                stage_graph(state, args)
            elif stage_name == "load":
                stage_load(state, args)
            elif stage_name == "enrichment":
                stage_enrichment(state, args)
            elif stage_name == "retrieval":
                stage_retrieval(state, args)
            elif stage_name == "visualization":
                stage_visualization(state, args)

        except Exception as e:
            logger.error(f"STAGE {stage_name} FAILED: {e}")
            if args.fail_fast:
                raise
            logger.info("Continuing to next stage...")

    # Final reports
    generate_final_reports(state, args)

    elapsed = time.time() - state.start_time
    logger.info("=" * 70)
    logger.info(f"  PIPELINE COMPLETE — {elapsed:.1f}s")
    logger.info(f"  Artifacts: {state.artifact_dir}")
    logger.info(f"  Failures: {len(state.failures)}")
    logger.info("=" * 70)


if __name__ == "__main__":
    import logging
    logging.basicConfig(
        level=logging.INFO,
        format="%(levelname)s %(name)s: %(message)s",
    )
    main()
