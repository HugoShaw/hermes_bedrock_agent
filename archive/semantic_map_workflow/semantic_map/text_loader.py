"""
Multi-format text loader for the semantic map workflow.

Dispatches to format-specific extraction logic based on file extension.
Optional heavy dependencies (pypdf/pdfplumber, python-docx, openpyxl,
python-pptx) are imported lazily so the module is still usable when those
packages are not installed.
"""

from __future__ import annotations

import csv
import io
import json
import logging
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Optional dependency availability flags
# ---------------------------------------------------------------------------
try:
    import pdfplumber  # type: ignore
    _PDF_LOADER = "pdfplumber"
except ImportError:
    pdfplumber = None  # type: ignore
    try:
        from pypdf import PdfReader as _PdfReader  # type: ignore
        _PDF_LOADER = "pypdf"
    except ImportError:
        _PdfReader = None  # type: ignore
        _PDF_LOADER = None

try:
    from docx import Document as _DocxDocument  # type: ignore
    _DOCX_AVAILABLE = True
except ImportError:
    _DocxDocument = None  # type: ignore
    _DOCX_AVAILABLE = False

try:
    import openpyxl  # type: ignore
    _XLSX_AVAILABLE = True
except ImportError:
    openpyxl = None  # type: ignore
    _XLSX_AVAILABLE = False

try:
    from pptx import Presentation as _PptxPresentation  # type: ignore
    _PPTX_AVAILABLE = True
except ImportError:
    _PptxPresentation = None  # type: ignore
    _PPTX_AVAILABLE = False


# ---------------------------------------------------------------------------
# Extension sets
# ---------------------------------------------------------------------------

# Extensions handled by plain-text reading
_PLAIN_TEXT_EXTENSIONS: frozenset[str] = frozenset(
    {
        ".txt", ".md", ".csv", ".json", ".jsonl",
        ".xml", ".yaml", ".yml", ".sql",
        ".java", ".py", ".js", ".ts", ".tsx", ".jsx",
        ".html", ".htm", ".properties", ".conf", ".ini",
        ".sh", ".bat",
    }
)


# ---------------------------------------------------------------------------
# Format-specific loaders
# ---------------------------------------------------------------------------

def _load_plain_text(path: str) -> str:
    """Read the file as UTF-8 text, falling back to latin-1."""
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return Path(path).read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    logger.warning("_load_plain_text: could not decode %s with any known encoding", path)
    return ""


def _load_csv(path: str) -> str:
    """Return CSV file content with headers preserved as plain text."""
    try:
        with open(path, newline="", encoding="utf-8") as fh:
            reader = csv.reader(fh)
            rows = list(reader)
        return "\n".join(",".join(row) for row in rows)
    except Exception as exc:
        logger.warning("_load_csv failed for %s: %s", path, exc)
        return ""


def _load_pdf(path: str) -> str:
    """Extract text from every page of a PDF."""
    if _PDF_LOADER == "pdfplumber":
        try:
            with pdfplumber.open(path) as pdf:
                pages = [page.extract_text() or "" for page in pdf.pages]
            return "\n\n".join(pages)
        except Exception as exc:
            logger.warning("pdfplumber failed for %s: %s", path, exc)
            return ""

    if _PDF_LOADER == "pypdf":
        try:
            reader = _PdfReader(path)
            pages = [page.extract_text() or "" for page in reader.pages]
            return "\n\n".join(pages)
        except Exception as exc:
            logger.warning("pypdf failed for %s: %s", path, exc)
            return ""

    logger.warning("load_text: no PDF library available; skipping %s", path)
    return ""


def _load_docx(path: str) -> str:
    """Extract paragraph text from a .docx file."""
    if not _DOCX_AVAILABLE:
        logger.warning("load_text: python-docx not installed; skipping %s", path)
        return ""
    try:
        doc = _DocxDocument(path)
        return "\n".join(para.text for para in doc.paragraphs)
    except Exception as exc:
        logger.warning("_load_docx failed for %s: %s", path, exc)
        return ""


def _load_xlsx(path: str) -> str:
    """Extract all sheets from an .xlsx file as plain text with column headers."""
    if not _XLSX_AVAILABLE:
        logger.warning("load_text: openpyxl not installed; skipping %s", path)
        return ""
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            headers = [str(cell) if cell is not None else "" for cell in rows[0]]
            part_lines: list[str] = [f"=== Sheet: {sheet_name} ==="]
            part_lines.append("\t".join(headers))
            for row in rows[1:]:
                part_lines.append(
                    "\t".join(str(cell) if cell is not None else "" for cell in row)
                )
            parts.append("\n".join(part_lines))
        wb.close()
        return "\n\n".join(parts)
    except Exception as exc:
        logger.warning("_load_xlsx failed for %s: %s", path, exc)
        return ""


def _load_pptx(path: str) -> str:
    """Extract text from every slide of a .pptx file."""
    if not _PPTX_AVAILABLE:
        logger.warning("load_text: python-pptx not installed; skipping %s", path)
        return ""
    try:
        prs = _PptxPresentation(path)
        slides: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            if texts:
                slides.append(f"--- Slide {idx} ---\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as exc:
        logger.warning("_load_pptx failed for %s: %s", path, exc)
        return ""


def _load_xml(path: str) -> str:
    """Return the raw XML text (the generic case is just plain text)."""
    return _load_plain_text(path)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def load_text(path: str) -> str:
    """Load the content of a file as a string, dispatching by extension.

    Supported formats:
      - Plain text: txt, md, csv, json, jsonl, xml, yaml, yml, sql,
        java, py, js, ts, tsx, jsx, html, htm, properties, conf, ini, sh, bat
      - PDF:  .pdf  (requires pdfplumber or pypdf)
      - Word: .docx (requires python-docx)
      - Excel: .xlsx, .xls (requires openpyxl)
      - PowerPoint: .pptx (requires python-pptx)

    Returns an empty string for unsupported formats or unreadable files.

    Parameters
    ----------
    path:
        Absolute or relative path to the file.

    Returns
    -------
    str
        Extracted text content.
    """
    p = Path(path)
    if not p.exists():
        logger.warning("load_text: file does not exist: %s", path)
        return ""

    ext = p.suffix.lower()

    if ext == ".csv":
        return _load_csv(path)
    if ext in _PLAIN_TEXT_EXTENSIONS:
        return _load_plain_text(path)
    if ext == ".pdf":
        return _load_pdf(path)
    if ext == ".docx":
        return _load_docx(path)
    if ext in {".xlsx", ".xls"}:
        return _load_xlsx(path)
    if ext == ".pptx":
        return _load_pptx(path)

    logger.warning("load_text: unsupported file format '%s'; skipping %s", ext, path)
    return ""


def load_chunks(
    path: str,
    max_chars: int = 4000,
    overlap: int = 200,
) -> list[str]:
    """Load a file and split its text into overlapping character-level chunks.

    Parameters
    ----------
    path:
        Path to the source file.
    max_chars:
        Maximum number of characters per chunk.
    overlap:
        Number of characters to repeat at the start of each successive chunk
        to preserve context across boundaries.

    Returns
    -------
    list[str]
        List of text chunks.  Returns ``[""]`` when the file is empty.
    """
    text = load_text(path)
    if not text:
        return []

    return _split_text(text, max_chars=max_chars, overlap=overlap)


def _split_text(text: str, max_chars: int = 4000, overlap: int = 200) -> list[str]:
    """Split *text* into chunks of at most *max_chars* characters with *overlap*."""
    if len(text) <= max_chars:
        return [text]

    overlap = min(overlap, max_chars // 2)
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = start + max_chars
        chunk = text[start:end]
        chunks.append(chunk)
        if end >= len(text):
            break
        start = end - overlap

    return chunks
